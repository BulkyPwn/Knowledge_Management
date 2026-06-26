"""
PPT Style Extractor — Extract visual style from a reference PPTX.

Extracts: theme colors, font families, slide dimensions, light/dark judgment.
Used to seed _generate_visual_design() with reference-derived values instead of LLM defaults.

Usage:
    from ppt_style_extractor import extract_pptx_style
    style = extract_pptx_style("/path/to/reference.pptx")
    # → {"colors": {...}, "fonts": {...}, "theme": "dark", ...}
"""

from __future__ import annotations

import json
import os
import re
import subprocess
import sys
from pathlib import Path
from typing import Dict, Any, Optional

def _default_ppt_master_dir() -> str:
    """Prefer the bundled ppt-master assets, fall back to the user-installed skill."""
    local_dir = os.path.join(os.path.dirname(__file__), "ppt_master")
    if os.path.isdir(local_dir):
        return local_dir
    return os.path.join(os.environ.get("APPDATA", ""), "chrys", "skills", "ppt-master")


SKILL_DIR = os.environ.get("PPT_MASTER_SKILL_DIR", _default_ppt_master_dir())


def extract_pptx_style(pptx_path: str) -> Dict[str, Any]:
    """
    Extract visual style from a reference PPTX file.

    Returns a dict compatible with _generate_visual_design()'s output format:
    {
        "theme": "dark" | "light",
        "colors": {"bg": "...", "primary": "...", ...},
        "typography": {"font_stack": "...", "title_stack": "...", "body_size": N},
        "icons": [...],
        "design_rationale": "从参考PPT提取的样式",
        "slide_count": N,
        "reference_path": "path/to/pptx",
    }
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return _fallback_style()

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"[style_extractor] Cannot open PPTX: {e}")
        return _fallback_style()

    result = {
        "extracted": True,
        "theme": _detect_theme(prs),
        "colors": _extract_colors(prs),
        "typography": _extract_typography(prs),
        "icons": _default_icons(),
        "design_rationale": "从参考PPT提取的样式",
        "slide_count": len(prs.slides),
        "reference_path": pptx_path,
    }
    return result


def extract_content(pptx_path: str) -> Optional[str]:
    """Extract text content from PPTX as Markdown using ppt_to_md.py."""
    script = os.path.join(SKILL_DIR, "scripts", "source_to_md", "ppt_to_md.py")
    if not os.path.exists(script):
        return None

    try:
        result = subprocess.run(
            [sys.executable, script, pptx_path, "--output", "-"],
            capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=120,
        )
        if result.returncode == 0:
            return result.stdout
        print(f"[style_extractor] ppt_to_md failed: {result.stderr[:200]}")
    except Exception as e:
        print(f"[style_extractor] ppt_to_md error: {e}")
    return None


def convert_to_svg(pptx_path: str, output_dir: str) -> bool:
    """Convert reference PPTX to per-slide SVG using pptx_to_svg.py."""
    script = os.path.join(SKILL_DIR, "scripts", "pptx_to_svg.py")
    if not os.path.exists(script):
        return False

    try:
        result = subprocess.run(
            [sys.executable, script, pptx_path, "-o", output_dir,
             "--inheritance-mode", "flat"],
            capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=300,
        )
        if result.returncode == 0:
            return True
        print(f"[style_extractor] pptx_to_svg failed: {result.stderr[:200]}")
    except Exception as e:
        print(f"[style_extractor] pptx_to_svg error: {e}")
    return False


# ---------------------------------------------------------------------------
# Internal extraction helpers
# ---------------------------------------------------------------------------

def _detect_theme(prs) -> str:
    """Detect if reference is dark or light theme by sampling slide backgrounds."""
    dark_count = 0
    light_count = 0
    sampled = 0

    for slide in prs.slides[:5]:  # sample first 5 slides
        bg = slide.background
        if bg.fill.type is not None:
            try:
                fill = bg.fill
                if hasattr(fill, 'fore_color') and fill.fore_color.type is not None:
                    rgb = fill.fore_color.rgb
                    brightness = (rgb[0] + rgb[1] + rgb[2]) / 3
                    if brightness < 128:
                        dark_count += 1
                    else:
                        light_count += 1
                    sampled += 1
            except Exception:
                pass

    if sampled == 0:
        return "dark"  # default
    return "dark" if dark_count >= light_count else "light"


def _extract_colors(prs) -> Dict[str, str]:
    """Extract theme colors from PPTX. Returns HEX dict matching visual_design format."""
    colors = _default_colors()

    try:
        # Try to read theme colors from slide master
        for slide_layout in prs.slide_masters[0].slide_layouts[:1] if prs.slide_masters else []:
            pass

        # Sample colors from shapes on the first few slides
        sampled_fills = []
        sampled_texts = []

        for slide in list(prs.slides)[:5]:
            for shape in slide.shapes:
                try:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    if run.font.color and run.font.color.type is not None:
                                        rgb = run.font.color.rgb
                                        sampled_texts.append(f"#{rgb}")
                                except Exception:
                                    pass
                    if hasattr(shape, 'fill'):
                        try:
                            if shape.fill.type is not None and hasattr(shape.fill, 'fore_color'):
                                rgb = shape.fill.fore_color.rgb
                                hex_color = f"#{rgb}"
                                brightness = (rgb[0] + rgb[1] + rgb[2]) / 3
                                sampled_fills.append((hex_color, brightness))
                        except Exception:
                            pass
                except Exception:
                    pass

        # Classify fills: darkest → bg, brightest → text, mid → accents
        if sampled_fills:
            sampled_fills.sort(key=lambda x: x[1])
            dark_palette = sampled_fills[:3]  # darkest
            light_palette = sampled_fills[-3:]  # brightest

            if dark_palette:
                colors["bg"] = dark_palette[0][0]
                if len(dark_palette) >= 2:
                    colors["secondary_bg"] = dark_palette[1][0]
            if light_palette:
                colors["text"] = light_palette[-1][0]
                if len(light_palette) >= 2:
                    colors["text_secondary"] = light_palette[-2][0]

        # Extract accent-like colors (mid-brightness, distinctive)
        if sampled_fills:
            mid = [c for c, b in sampled_fills if 60 < b < 200]
            if len(mid) >= 2:
                colors["primary"] = mid[0]
                if len(mid) >= 3:
                    colors["accent"] = mid[1]
                    colors["secondary_accent"] = mid[2]
                elif len(mid) >= 2:
                    colors["accent"] = mid[1]

        # Handle text colors
        if sampled_texts and colors["text"] == _default_colors()["text"]:
            colors["text"] = sampled_texts[-1]
            if len(sampled_texts) >= 2:
                colors["text_secondary"] = sampled_texts[0]

    except Exception as e:
        print(f"[style_extractor] Color extraction fallback: {e}")

    return colors


def _extract_typography(prs) -> Dict[str, Any]:
    """Extract font families and sizes from reference PPTX."""
    fonts = {"font_stack": "", "title_stack": "", "body_size": 18, "title_size": 32}
    font_counts = {}
    size_values = []

    try:
        for slide in list(prs.slides)[:10]:
            for shape in slide.shapes:
                try:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    if run.font.name:
                                        font_counts[run.font.name] = font_counts.get(run.font.name, 0) + 1
                                    if run.font.size:
                                        size_pt = run.font.size / 12700  # EMU → pt
                                        size_px = int(size_pt * 1.333)  # pt → px (approx)
                                        size_values.append(size_px)
                                except Exception:
                                    pass
                except Exception:
                    pass
    except Exception:
        pass

    # Get top fonts
    sorted_fonts = sorted(font_counts.items(), key=lambda x: x[1], reverse=True)
    if sorted_fonts:
        top_font = sorted_fonts[0][0]
        fonts["font_stack"] = f'"{top_font}", "Microsoft YaHei", Arial, sans-serif'
        # Title font: use second most common, or Georgia for serif contrast
        if len(sorted_fonts) >= 2:
            second_font = sorted_fonts[1][0]
            fonts["title_stack"] = f'"{second_font}", "Microsoft YaHei", serif'
        else:
            fonts["title_stack"] = f'Georgia, "{top_font}", serif'

    # Get median size as body baseline
    if size_values:
        size_values.sort()
        fonts["body_size"] = size_values[len(size_values) // 2]
        fonts["title_size"] = max(28, int(fonts["body_size"] * 1.75))

    return fonts


def _default_colors() -> Dict[str, str]:
    return {
        "bg": "#0B1120", "secondary_bg": "#131E33",
        "primary": "#3B82F6", "accent": "#06B6D4", "secondary_accent": "#8B5CF6",
        "text": "#E2E8F0", "text_secondary": "#94A3B8", "border": "#1E3A5F",
        "success": "#10B981", "warning": "#EF4444",
    }


def _default_icons() -> list:
    return ["bulb", "chart-bar", "award", "code", "server", "apps", "shield",
            "rocket", "target", "list", "star", "check", "bolt", "settings",
            "database", "graph", "globe", "user", "link", "search"]


def _fallback_style() -> Dict[str, Any]:
    return {
        "extracted": False,
        "theme": "dark",
        "colors": _default_colors(),
        "typography": {
            "font_stack": '"Microsoft YaHei", "PingFang SC", Arial, sans-serif',
            "title_stack": 'Georgia, "Microsoft YaHei", serif',
            "body_size": 18, "title_size": 32,
        },
        "icons": _default_icons(),
        "design_rationale": "默认深色科技主题",
        "slide_count": 0,
        "reference_path": "",
    }
