"""
Image to Mermaid diagram converter.

Scans raw/assets/ for image files, uses Vision LLM to identify
diagrams (logic view, sequence diagram, flowchart) and generates
corresponding Mermaid code files.
"""

import base64
import json
import os
import re
import shutil
import subprocess
import tempfile
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from typing import Any, Dict, List, Optional

import requests


# ── Thinking 限制配置 ──────────────────────────────────────────
_DEFAULT_THINKING_LIMITS = {
    "vision_classify": 2048,
    "vision_mermaid": 4096,
    "vision_describe": 4096,
    "vision_convert": 4096,
    "chat_qa": 0,
    "ppt_generation": 0,
    "search_research": 0,
    "agent_loop": 0,
}

def get_thinking_limit(scene_key: str, default: int = 0) -> int:
    """从 ~/.SSSC_AI/knowledge_management.json 读取 thinking 限制配置。

    支持两种格式：
    - 嵌套: {"thinking_limits": {"vision_classify": {"max_tokens": 2048}}}
    - 简洁: {"thinking_limits": {"vision_classify": 2048}}

    Returns: max_tokens 值（0 = 不限制）
    """
    try:
        config_dir = os.path.join(os.path.expanduser("~"), ".SSSC_AI")
        config_path = os.path.join(config_dir, "knowledge_management.json")
        if not os.path.exists(config_path):
            return _DEFAULT_THINKING_LIMITS.get(scene_key, default)
        with open(config_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        limits = data.get("thinking_limits", {})
        if scene_key not in limits:
            return _DEFAULT_THINKING_LIMITS.get(scene_key, default)
        val = limits[scene_key]
        # 支持嵌套格式 {"max_tokens": N} 和简洁格式 N
        if isinstance(val, dict):
            return int(val.get("max_tokens", _DEFAULT_THINKING_LIMITS.get(scene_key, default)))
        return int(val)
    except Exception:
        return _DEFAULT_THINKING_LIMITS.get(scene_key, default)


# Supported image extensions
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"}

# Mermaid diagram types we care about
MERMAID_DIAGRAM_TYPES = [
    "flowchart", "graph",
    "sequenceDiagram",
    "classDiagram",
    "stateDiagram", "stateDiagram-v2",
    "erDiagram",
    "journey",
    "gantt",
    "pie",
    "gitGraph",
    "mindmap",
    "timeline",
    "quadrantChart",
    "xychart-beta",
    "C4Context", "C4Container", "C4Component", "C4Dynamic", "C4Deployment",
]

# System prompt for the Vision LLM
SYSTEM_PROMPT = """你是一个专业的图表识别和结构化分析专家。

## 任务
分析图片，判断它是否为可转换的结构化图表。如果是，输出**结构化 JSON 数据**（而非 Mermaid 代码），系统将自动生成 Mermaid 代码。

## 支持的图表类型
- flowchart：流程图
- sequenceDiagram：时序图
- classDiagram：类图
- stateDiagram：状态图
- erDiagram：ER 图（实体关系图）
- journey：旅程图
- mindmap：思维导图
- timeline：时间线图

## 不是图表的情况（返回 is_diagram: false）
- 照片、截图、UI 界面截图
- 纯文字段落、表格
- 数据可视化图表（柱状图、饼图、折线图等统计图表）
- 手绘涂鸦、装饰性图案
- 无法清晰辨认结构和节点的图片

## 输出格式
严格按以下 JSON 格式返回，不要添加任何额外文字。

### 不是图表时：
```json
{
  "is_diagram": false,
  "diagram_type": "",
  "description": "图片内容简要描述",
  "diagram_data": null
}
```

### 是图表时，diagram_data 的结构取决于 diagram_type：

#### flowchart（流程图）
```json
{
  "is_diagram": true,
  "diagram_type": "flowchart",
  "description": "一句话描述",
  "diagram_data": {
    "direction": "TD",
    "nodes": [
      {"id": "A", "label": "节点文本", "shape": "rect"}
    ],
    "edges": [
      {"from": "A", "to": "B", "label": "连线标注", "style": "solid"}
    ],
    "subgraphs": [
      {"id": "sg1", "label": "子图名称", "nodes": ["A", "B"]}
    ]
  }
}
```
字段说明：
- direction: "TD"(从上到下)、"LR"(从左到右)、"BT"(从下到上)、"RL"(从右到左)
- nodes[].id: 只使用英文字母、数字和下划线（如 A、B1、node_2）
- nodes[].shape: "rect"(矩形)、"rounded"(圆角)、"stadium"(体育场)、"circle"(圆形)、"diamond"(菱形)、"hexagon"(六边形)、"parallelogram"(平行四边形)
- edges[].style: "solid"(实线箭头 -->)、"dotted"(虚线箭头 -.->)、"thick"(粗线箭头 ==>)
- edges[].label: 可选，连线上的标注文字
- subgraphs: 可选，原图有明显分组时使用
- 中文标签和标注保持原样，不要翻译

#### sequenceDiagram（时序图）
```json
{
  "is_diagram": true,
  "diagram_type": "sequenceDiagram",
  "description": "一句话描述",
  "diagram_data": {
    "participants": [
      {"id": "C", "label": "客户端"}
    ],
    "messages": [
      {"from": "C", "to": "S", "text": "请求数据", "type": "solid_arrow"}
    ]
  }
}
```
- messages[].type: "solid_arrow"(实线 ->>), "dashed_arrow"(虚线 -->>), "solid_open"(实线开放箭头 ->), "dashed_open"(虚线开放箭头 -->)
- 可选: messages 中可包含 "note": {"position": "right_of"|"left_of"|"over", "participant": "C", "text": "备注内容"}
- 可选: diagram_data 中可包含 "groups": [{"type": "loop"|"alt"|"opt"|"par", "label": "条件", "messages": [...]}]

#### classDiagram（类图）
```json
{
  "is_diagram": true,
  "diagram_type": "classDiagram",
  "description": "一句话描述",
  "diagram_data": {
    "classes": [
      {"name": "ClassName", "members": ["+field: Type", "-private_method()"], "annotations": ["abstract"]}
    ],
    "relationships": [
      {"from": "ClassA", "to": "ClassB", "type": "inheritance", "label": ""}
    ]
  }
}
```
- relationships[].type: "inheritance"(<|--), "composition"(*--), "aggregation"(o--), "association"(-->), "dependency"(..>), "realization"(..|>)

#### stateDiagram（状态图）
```json
{
  "is_diagram": true,
  "diagram_type": "stateDiagram",
  "description": "一句话描述",
  "diagram_data": {
    "states": [
      {"id": "s1", "label": "状态1"}
    ],
    "transitions": [
      {"from": "[*]", "to": "s1", "label": ""},
      {"from": "s1", "to": "[*]", "label": "完成"}
    ]
  }
}
```
- "[*]" 表示起始/结束状态

#### erDiagram（ER 图）
```json
{
  "is_diagram": true,
  "diagram_type": "erDiagram",
  "description": "一句话描述",
  "diagram_data": {
    "entities": [
      {"name": "USER", "attributes": [{"type": "string", "name": "username", "key": "PK"}]}
    ],
    "relationships": [
      {"from": "USER", "to": "ORDER", "from_cardinality": "one", "to_cardinality": "many", "label": "places"}
    ]
  }
}
```
- cardinality: "one"、"many"、"one_or_many"、"zero_or_one"、"zero_or_many"

#### mindmap（思维导图）
```json
{
  "is_diagram": true,
  "diagram_type": "mindmap",
  "description": "一句话描述",
  "diagram_data": {
    "root": {
      "label": "中心主题",
      "children": [
        {"label": "分支1", "children": [{"label": "子节点"}]},
        {"label": "分支2", "children": []}
      ]
    }
  }
}
```

#### journey（旅程图）
```json
{
  "is_diagram": true,
  "diagram_type": "journey",
  "description": "一句话描述",
  "diagram_data": {
    "title": "用户旅程",
    "sections": [
      {"name": "阶段1", "tasks": [{"name": "任务A", "score": 5, "actors": ["用户"]}]}
    ]
  }
}
```
- score: 1-5 的满意度评分

#### timeline（时间线图）
```json
{
  "is_diagram": true,
  "diagram_type": "timeline",
  "description": "一句话描述",
  "diagram_data": {
    "title": "项目时间线",
    "periods": [
      {"label": "2024 Q1", "events": ["事件1", "事件2"]}
    ]
  }
}
```

## 重要规则
1. 节点 ID / name 只使用英文字母、数字和下划线
2. 中文内容放在 label / text / description 等文本字段中
3. 包含图中所有可见的节点和连接关系，不要遗漏
4. diagram_data 中的字段如果原图没有，可以省略（如 subgraphs、groups）"""

USER_PROMPT = "请分析这张图片，判断是否为图表。如果是，请输出结构化 JSON 数据。"

# ============================================================================
#  Three-step strategy prompts
# ============================================================================

# Step 1: 复杂度分类 Prompt
CLASSIFY_SYSTEM_PROMPT = """你是一个专业的图表识别和分类专家。

## 任务
分析图片，判断它是否为可转换的结构化图表，并评估其复杂度。

## 支持的图表类型
flowchart（流程图）、sequenceDiagram（时序图）、classDiagram（类图）、stateDiagram（状态图）、erDiagram（ER 图）、journey（旅程图）、mindmap（思维导图）、timeline（时间线图）

## 不是图表的情况
- 照片、截图、UI 界面截图
- 纯文字段落、表格
- 数据可视化图表（柱状图、饼图、折线图等统计图表）
- 手绘涂鸦、装饰性图案
- 无法清晰辨认结构和节点的图片

## 复杂度判断标准
**simple（简单）**:
- 节点数量 ≤ 8
- 没有子图（subgraph）
- 没有复杂的分组或嵌套结构
- 连线关系清晰简单
- 节点形状种类 ≤ 2

**complex（复杂）**:
- 节点数量 > 8
- 包含子图（subgraph）或分组
- 有多种不同的节点形状
- 包含嵌套结构
- 连线关系复杂（有多种样式、条件分支）
- 时序图中有多个 alt/loop/opt/par 分组

## 输出格式
严格按以下 JSON 格式返回，不要添加任何额外文字：
```json
{
  "is_diagram": true/false,
  "complexity": "simple" 或 "complex",
  "diagram_type": "图表类型",
  "description": "一句话描述图片内容",
  "node_count_estimate": 大致节点数
}
```
如果不是图表，is_diagram 为 false，complexity 为空字符串。"""

CLASSIFY_USER_PROMPT = "请分析这张图片，判断是否为图表并评估复杂度。"

# Step 2 Plan B: 直接生成 Mermaid 代码 Prompt
DIRECT_MERMAID_SYSTEM_PROMPT = """你是一个专业的 Mermaid 图表代码生成专家。

## 任务
分析图片中的图表，直接输出有效的 Mermaid 代码。

## 输出规则
1. 只输出一个 ```mermaid ... ``` 代码块，不要有任何其他文字
2. 代码必须语法正确，可被 Mermaid 渲染引擎直接渲染
3. 节点 ID 只使用英文字母、数字和下划线（如 A、B1、node_2）
4. **禁止使用 Mermaid 保留关键字作为节点 ID**：start, end, subgraph, if, else, loop, alt, opt, par, and, or, not
5. 包含特殊字符的标签用双引号包裹，如 A["标签(含特殊字符)"]
6. 中文内容保留原样
7. 包含图中所有可见的节点和连接关系，不要遗漏

## 各类型图表格式

### flowchart
```mermaid
flowchart TD
    A[开始] --> B{判断}
    B -->|是| C[执行]
    B -->|否| D[结束]
```
- 方向: TD/LR/BT/RL
- 节点形状: []矩形、()圆角、([])体育场、(())圆形、{}菱形、{{}}六边形、[/]平行四边形
- 连线: -->实线、-.->虚线、==>粗线
- 子图用 subgraph name["标签"] ... end

### sequenceDiagram
```mermaid
sequenceDiagram
    participant C as 客户端
    participant S as 服务端
    C->>S: 请求数据
    S-->>C: 返回数据
```
- 箭头: ->>实线、-->>虚线、->开放、-->虚线开放
- 分组: loop/alt/opt/par ... end
- 备注: Note right of/left of/over P: 文本

### classDiagram
```mermaid
classDiagram
    class Animal {
        +String name
        +eat()
    }
    Animal <|-- Dog
```
- 关系: <|--继承、*--组合、o--聚合、-->关联、..>依赖、..|>实现

### stateDiagram-v2
```mermaid
stateDiagram-v2
    [*] --> Idle
    Idle --> Processing : start
    Processing --> [*] : done
```

### erDiagram
```mermaid
erDiagram
    USER {
        string username PK
        string email
    }
    USER ||--o{ ORDER : places
```

### mindmap
```mermaid
mindmap
  root((中心主题))
    分支1
      子节点A
    分支2
      子节点B
```

### journey
```mermaid
journey
    title 用户旅程
    section 阶段1
      任务A: 5: 用户
```

### timeline
```mermaid
timeline
    title 项目时间线
    2024 Q1
        : 事件1
        : 事件2
```

## 重要
- 只输出 ```mermaid ... ``` 代码块
- 不要输出任何解释、注释或其他文字"""

DIRECT_MERMAID_USER_PROMPT = "请分析这张图片中的图表，直接输出有效的 Mermaid 代码。"

# Plan B 纠错 Prompt: LLM 校验失败后，把报错 + 错误代码发给 LLM 修正
FIX_MERMAID_SYSTEM_PROMPT = """你是一个专业的 Mermaid 图表代码修复专家。

## 任务
上一次生成的 Mermaid 代码有语法错误，无法被渲染引擎解析。请根据图片内容和错误信息，修复代码。

## 输出规则
1. 只输出一个 ```mermaid ... ``` 代码块，不要有任何其他文字
2. 修复后的代码必须语法正确
3. 节点 ID 只使用英文字母、数字和下划线
4. **禁止使用 Mermaid 保留关键字作为节点 ID**：start, end, subgraph, if, else, loop, alt, opt, par, and, or, not
5. 包含特殊字符的标签用双引号包裹
6. 中文内容保留原样
7. 保持原图的完整结构，不要删减节点和连线

## 常见语法错误及修复方法
- 箭头不完整: `-->` 后面必须有目标节点
- 缺少 end: subgraph、loop、alt、opt、par 都必须有对应的 end
- 括号不匹配: `{` 和 `}` 数量必须一致
- 空标签: 标签不能为空字符串
- 特殊字符未转义: 括号、冒号等在标签中需要双引号包裹
- 节点 ID 包含非法字符: 只允许字母、数字、下划线
- **使用保留关键字作为节点 ID**: start/end 等是 Mermaid 保留字，需改用 startNode/endNode 等

## 重要
- 只输出 ```mermaid ... ``` 代码块
- 不要输出任何解释、注释或其他文字"""


def _build_fix_user_prompt(error_msg: str, bad_code: str) -> str:
    """构建纠错用户提示词，包含错误信息和错误代码。"""
    return (
        f"上一次生成的 Mermaid 代码有语法错误，请修正。\n\n"
        f"## 错误信息\n{error_msg}\n\n"
        f"## 有错误的代码\n```mermaid\n{bad_code}\n```\n\n"
        f"请对照图片，修复以上代码并输出正确的 Mermaid 代码。"
    )


# Step 3 Fallback: 简化 JSON Prompt（兜底用）
FALLBACK_JSON_USER_PROMPT = "这张图片包含一个图表。请用最简化的结构化 JSON 输出（减少节点数量，保留核心结构）。"


# ============================================================================
#  Structured JSON → Mermaid code converters
# ============================================================================

def _escape_label(text: str) -> str:
    """Wrap label in double quotes if it contains special characters."""
    if not text:
        return '""'
    needs_quote = any(c in text for c in '()[]{}:;#&|<>"')
    if needs_quote:
        escaped = text.replace('"', '\\"')
        return f'"{escaped}"'
    return text


def _flowchart_shape_to_syntax(shape: str) -> tuple:
    """Return (open, close) brackets for a flowchart node shape."""
    mapping = {
        "rect": ("[", "]"),
        "rounded": ("(", ")"),
        "stadium": ("([", "])"),
        "circle": ("((", "))"),
        "diamond": ("{", "}"),
        "hexagon": ("{{", "}}"),
        "parallelogram": ("[/", "/]"),
    }
    return mapping.get(shape, ("[", "]"))


def _flowchart_edge_syntax(style: str, has_label: bool, label: str) -> str:
    """Return the edge arrow syntax for flowchart."""
    arrows = {
        "solid": "-->",
        "dotted": "-.->",
        "thick": "==>",
    }
    arrow = arrows.get(style, "-->")
    if has_label and label:
        return f"{arrow}|{_escape_label(label)}|"
    return arrow


def _gen_flowchart(data: dict) -> str:
    """Generate Mermaid flowchart code from structured data."""
    lines = []
    direction = data.get("direction", "TD").upper()
    if direction not in ("TD", "LR", "BT", "RL"):
        direction = "TD"
    lines.append(f"flowchart {direction}")

    nodes = data.get("nodes", [])
    edges = data.get("edges", [])
    subgraphs = data.get("subgraphs", [])

    # Build a map of node id → definition line
    node_defs = {}
    for node in nodes:
        nid = node.get("id", "")
        if not nid:
            continue
        label = node.get("label", nid)
        shape = node.get("shape", "rect")
        open_b, close_b = _flowchart_shape_to_syntax(shape)
        node_defs[nid] = f"    {nid}{open_b}{_escape_label(label)}{close_b}"

    # Build subgraph → node set map
    subgraph_nodes = set()
    for sg in subgraphs:
        for nid in sg.get("nodes", []):
            subgraph_nodes.add(nid)

    # Emit subgraphs first
    for sg in subgraphs:
        sg_id = sg.get("id", "sg")
        sg_label = sg.get("label", "")
        lines.append(f"    subgraph {sg_id}[{_escape_label(sg_label)}]")
        for nid in sg.get("nodes", []):
            if nid in node_defs:
                lines.append(node_defs[nid])
            else:
                lines.append(f"        {nid}")
        lines.append("    end")

    # Emit standalone nodes (not in any subgraph)
    for nid, definition in node_defs.items():
        if nid not in subgraph_nodes:
            lines.append(definition)

    # Emit edges
    for edge in edges:
        src = edge.get("from", "")
        dst = edge.get("to", "")
        if not src or not dst:
            continue
        label = edge.get("label", "")
        style = edge.get("style", "solid")
        arrow = _flowchart_edge_syntax(style, bool(label), label)
        lines.append(f"    {src} {arrow} {dst}")

    return "\n".join(lines)


def _gen_sequence(data: dict) -> str:
    """Generate Mermaid sequenceDiagram code from structured data."""
    lines = ["sequenceDiagram"]

    for p in data.get("participants", []):
        pid = p.get("id", "")
        label = p.get("label", pid)
        if label and label != pid:
            lines.append(f"    participant {pid} as {_escape_label(label)}")
        else:
            lines.append(f"    participant {pid}")

    def _emit_messages(msgs: list, indent: str = "    "):
        for msg in msgs:
            # Handle notes
            note = msg.get("note")
            if note:
                pos = note.get("position", "right_of")
                participant = note.get("participant", "")
                text = note.get("text", "")
                pos_map = {"right_of": "right of", "left_of": "left of", "over": "over"}
                lines.append(f"{indent}Note {pos_map.get(pos, 'right of')} {participant}: {text}")
                continue

            src = msg.get("from", "")
            dst = msg.get("to", "")
            text = msg.get("text", "")
            msg_type = msg.get("type", "solid_arrow")
            arrow_map = {
                "solid_arrow": "->>",
                "dashed_arrow": "-->>",
                "solid_open": "->",
                "dashed_open": "-->",
            }
            arrow = arrow_map.get(msg_type, "->>")
            lines.append(f"{indent}{src}{arrow}{dst}: {text}")

    # Handle groups (loop, alt, opt, par)
    groups = data.get("groups", [])
    messages = data.get("messages", [])

    if groups:
        for group in groups:
            g_type = group.get("type", "loop")
            g_label = group.get("label", "")
            lines.append(f"    {g_type} {_escape_label(g_label)}")
            _emit_messages(group.get("messages", []), indent="    ")
            lines.append("    end")
    else:
        _emit_messages(messages)

    return "\n".join(lines)


def _gen_class(data: dict) -> str:
    """Generate Mermaid classDiagram code from structured data."""
    lines = ["classDiagram"]

    for cls in data.get("classes", []):
        name = cls.get("name", "")
        if not name:
            continue
        annotations = cls.get("annotations", [])
        for ann in annotations:
            lines.append(f"    class {name} {{")
            lines.append(f"        <<{ann}>>")
            break
        else:
            lines.append(f"    class {name} {{")

        for member in cls.get("members", []):
            lines.append(f"        {member}")
        lines.append("    }")

    rel_map = {
        "inheritance": "<|--",
        "composition": "*--",
        "aggregation": "o--",
        "association": "-->",
        "dependency": "..>",
        "realization": "..|>",
    }
    for rel in data.get("relationships", []):
        src = rel.get("from", "")
        dst = rel.get("to", "")
        rel_type = rel.get("type", "association")
        label = rel.get("label", "")
        arrow = rel_map.get(rel_type, "-->")
        line = f"    {src} {arrow} {dst}"
        if label:
            line += f" : {label}"
        lines.append(line)

    return "\n".join(lines)


def _gen_state(data: dict) -> str:
    """Generate Mermaid stateDiagram-v2 code from structured data."""
    lines = ["stateDiagram-v2"]

    for state in data.get("states", []):
        sid = state.get("id", "")
        label = state.get("label", "")
        if label and label != sid:
            lines.append(f"    {sid} : {label}")

    for tr in data.get("transitions", []):
        src = tr.get("from", "[*]")
        dst = tr.get("to", "[*]")
        label = tr.get("label", "")
        if label:
            lines.append(f"    {src} --> {dst} : {label}")
        else:
            lines.append(f"    {src} --> {dst}")

    return "\n".join(lines)


def _gen_er(data: dict) -> str:
    """Generate Mermaid erDiagram code from structured data."""
    lines = ["erDiagram"]

    for entity in data.get("entities", []):
        name = entity.get("name", "")
        if not name:
            continue
        attrs = entity.get("attributes", [])
        if attrs:
            lines.append(f"    {name} {{")
            for attr in attrs:
                atype = attr.get("type", "string")
                aname = attr.get("name", "")
                akey = attr.get("key", "")
                line = f"        {atype} {aname}"
                if akey:
                    line += f" {akey}"
                lines.append(line)
            lines.append("    }")
        else:
            lines.append(f"    {name}")

    card_map = {
        "one": "||",
        "many": "}o",
        "one_or_many": "|{",
        "zero_or_one": "o|",
        "zero_or_many": "o{",
    }
    # Reverse direction for "from" side
    card_from_map = {
        "one": "||",
        "many": "{o",
        "one_or_many": "{|",
        "zero_or_one": "|o",
        "zero_or_many": "{o",
    }

    for rel in data.get("relationships", []):
        src = rel.get("from", "")
        dst = rel.get("to", "")
        from_card = rel.get("from_cardinality", "one")
        to_card = rel.get("to_cardinality", "many")
        label = rel.get("label", "")
        left = card_from_map.get(from_card, "||")
        right = card_map.get(to_card, "}o")
        lines.append(f"    {src} {left}--{right} {dst} : {label}")

    return "\n".join(lines)


def _gen_mindmap(data: dict) -> str:
    """Generate Mermaid mindmap code from structured data."""
    lines = ["mindmap"]

    def _walk(node: dict, depth: int):
        label = node.get("label", "")
        indent = "    " + "  " * depth
        lines.append(f"{indent}{_escape_label(label)}")
        for child in node.get("children", []):
            _walk(child, depth + 1)

    root = data.get("root", {})
    if root:
        _walk(root, 0)

    return "\n".join(lines)


def _gen_journey(data: dict) -> str:
    """Generate Mermaid journey code from structured data."""
    lines = ["journey"]
    title = data.get("title", "")
    if title:
        lines.append(f"    title {title}")

    for section in data.get("sections", []):
        sec_name = section.get("name", "")
        lines.append(f"    section {sec_name}")
        for task in section.get("tasks", []):
            name = task.get("name", "")
            score = task.get("score", 3)
            actors = task.get("actors", [])
            task_line = f"        {name}: {score}"
            if actors:
                task_line += f": {', '.join(actors)}"
            lines.append(task_line)

    return "\n".join(lines)


def _gen_timeline(data: dict) -> str:
    """Generate Mermaid timeline code from structured data."""
    lines = ["timeline"]
    title = data.get("title", "")
    if title:
        lines.append(f"    title {title}")

    for period in data.get("periods", []):
        label = period.get("label", "")
        lines.append(f"        {label}")
        for event in period.get("events", []):
            lines.append(f"            : {event}")

    return "\n".join(lines)


# Dispatch table for diagram type → generator function
_DIAGRAM_GENERATORS = {
    "flowchart": _gen_flowchart,
    "graph": _gen_flowchart,
    "sequenceDiagram": _gen_sequence,
    "classDiagram": _gen_class,
    "stateDiagram": _gen_state,
    "stateDiagram-v2": _gen_state,
    "erDiagram": _gen_er,
    "mindmap": _gen_mindmap,
    "journey": _gen_journey,
    "timeline": _gen_timeline,
}


def _json_to_mermaid(llm_result: dict) -> Optional[str]:
    """Convert structured LLM response to Mermaid code.

    Returns Mermaid code string, or None if not a diagram or data is invalid.
    """
    if not llm_result:
        return None

    if not llm_result.get("is_diagram", False):
        return None

    diagram_type = llm_result.get("diagram_type", "")
    diagram_data = llm_result.get("diagram_data")

    if not diagram_type or not diagram_data:
        return None

    generator = _DIAGRAM_GENERATORS.get(diagram_type)
    if not generator:
        print(f"[image_to_desc] Unsupported diagram type: {diagram_type}")
        return None

    try:
        code = generator(diagram_data)
        if not code or len(code.strip()) < 10:
            print(f"[image_to_desc] Generated code too short for {diagram_type}")
            return None
        return code
    except Exception as e:
        print(f"[image_to_desc] Failed to generate Mermaid for {diagram_type}: {e}")
        return None


def _encode_image_to_data_url(image_path: str) -> Optional[str]:
    """Encode an image file to a base64 data URL."""
    try:
        ext = os.path.splitext(image_path)[1].lower()
        mime_map = {
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".webp": "image/webp",
            ".bmp": "image/bmp",
            ".tiff": "image/tiff",
        }
        mime = mime_map.get(ext, "image/png")
        with open(image_path, "rb") as f:
            data = f.read()
        b64 = base64.b64encode(data).decode("utf-8")
        return f"data:{mime};base64,{b64}"
    except Exception as e:
        print(f"[image_to_desc] Failed to encode {image_path}: {e}")
        return None


def _find_vision_model(models: List[dict]) -> Optional[dict]:
    """Find a model with vision=true from the models list."""
    for m in models:
        if m.get("vision", False):
            return m
    return None


def _build_chat_url(base_url: str) -> str:
    """Build chat/completions URL."""
    url = (base_url or "").strip().rstrip("/")
    if not url:
        return ""
    suffix = "/chat/completions"
    if url.endswith(suffix):
        url = url[:-len(suffix)]
    lower = url.lower()
    if lower.endswith("/api/coding/v3") or lower.endswith("/v3"):
        return f"{url}{suffix}"
    if not lower.endswith("/v1"):
        url = f"{url}/v1"
    return f"{url}{suffix}"


def _parse_llm_json_response(content: str) -> Optional[dict]:
    """从 LLM 返回的文本中提取 JSON 对象。

    解析顺序:
    1. 直接 json.loads
    2. ```json ... ``` 代码块
    3. 平衡花括号提取首个完整 JSON 对象
    """
    content = content.strip()

    # Try direct JSON parse
    try:
        return json.loads(content)
    except json.JSONDecodeError:
        pass

    # Try to extract JSON block from the response (```json ... ``` or raw { ... })
    json_block_match = re.search(r'```(?:json)?\s*(\{[\s\S]*?\})\s*```', content)
    if json_block_match:
        try:
            return json.loads(json_block_match.group(1))
        except json.JSONDecodeError:
            pass

    # Fallback: find first complete JSON object using balanced braces
    brace_depth = 0
    json_start = -1
    for i, c in enumerate(content):
        if c == '{':
            if brace_depth == 0:
                json_start = i
            brace_depth += 1
        elif c == '}':
            brace_depth -= 1
            if brace_depth == 0 and json_start >= 0:
                try:
                    return json.loads(content[json_start:i + 1])
                except json.JSONDecodeError:
                    json_start = -1

    print(f"[image_to_desc] Failed to parse LLM response as JSON, content_length={len(content)}, "
          f"content[:500]={content[:500]}", flush=True)
    return None


def _extract_mermaid_from_text(text: str) -> Optional[str]:
    """从 LLM 返回的文本中提取 ```mermaid ... ``` 代码块。"""
    text = text.strip()

    # 尝试匹配 ```mermaid ... ``` 代码块
    match = re.search(r'```mermaid\s*\n([\s\S]*?)```', text)
    if match:
        return match.group(1).strip()

    # 尝试匹配普通 ``` ... ``` 代码块（不含语言标记）
    match = re.search(r'```\s*\n([\s\S]*?)```', text)
    if match:
        content = match.group(1).strip()
        # 检查内容是否看起来像 mermaid
        if any(content.startswith(kw) for kw in [
            'flowchart', 'graph ', 'sequenceDiagram', 'classDiagram',
            'stateDiagram', 'erDiagram', 'mindmap', 'journey', 'timeline',
            'gantt', 'pie', 'gitGraph', 'quadrantChart', 'xychart',
            'C4Context', 'C4Container', 'C4Component', 'C4Dynamic', 'C4Deployment',
        ]):
            return content

    print(f"[image_to_desc] _extract_mermaid_from_text: 未找到 mermaid 代码块, "
          f"text[:500]={text[:500]}", flush=True)
    return None


# 缓存校验脚本路径和 node 可执行文件路径
_MERMAID_VALIDATE_SCRIPT: Optional[str] = None
_NODE_EXECUTABLE: Optional[str] = None


def _find_validate_script() -> Optional[str]:
    """查找 scripts/validate_mermaid.mjs 的路径。"""
    global _MERMAID_VALIDATE_SCRIPT
    if _MERMAID_VALIDATE_SCRIPT and os.path.isfile(_MERMAID_VALIDATE_SCRIPT):
        return _MERMAID_VALIDATE_SCRIPT

    # 从当前文件向上查找项目根目录
    probe = os.path.dirname(os.path.abspath(__file__))
    for _ in range(6):
        script_path = os.path.join(probe, "scripts", "validate_mermaid.mjs")
        if os.path.isfile(script_path):
            _MERMAID_VALIDATE_SCRIPT = script_path
            return script_path
        parent = os.path.dirname(probe)
        if parent == probe:
            break
        probe = parent

    return None


def _find_node() -> Optional[str]:
    """查找 node 可执行文件。"""
    global _NODE_EXECUTABLE
    if _NODE_EXECUTABLE:
        return _NODE_EXECUTABLE

    # 1. PATH 中的 node
    node = shutil.which("node")
    if node:
        _NODE_EXECUTABLE = node
        return node

    # 2. 项目自带的 Node.js (C:\nodejs\...)
    for candidate in [
        r"C:\nodejs\node-v22.16.0-win-x64\node.exe",
        r"C:\nodejs\node.exe",
    ]:
        if os.path.isfile(candidate):
            _NODE_EXECUTABLE = candidate
            return candidate

    return None


def _validate_mermaid(code: str) -> tuple:
    """使用 mermaid.parse() CLI 校验 Mermaid 语法。

    通过 `node scripts/validate_mermaid.mjs` 调用 mermaid 官方解析器，
    比手写正则校验更准确，能捕获所有语法错误。

    返回 (is_valid: bool, error_msg: str)。
    """
    if not code or len(code.strip()) < 15:
        return False, "代码太短"

    # 查找 node 和校验脚本
    node_exe = _find_node()
    if not node_exe:
        print("[image_to_desc] _validate_mermaid: node 未找到, 跳过校验", flush=True)
        return True, ""

    script_path = _find_validate_script()
    if not script_path:
        print("[image_to_desc] _validate_mermaid: validate_mermaid.mjs 未找到, 跳过校验", flush=True)
        return True, ""

    try:
        result = subprocess.run(
            [node_exe, script_path],
            input=code.encode("utf-8"),
            capture_output=True,
            timeout=15,
            cwd=os.path.dirname(os.path.dirname(script_path)),  # 项目根目录
        )

        if result.returncode != 0:
            stderr = result.stderr.decode("utf-8", errors="replace")[:200]
            print(f"[image_to_desc] _validate_mermaid: node 进程异常 exit={result.returncode}, stderr={stderr}", flush=True)
            return False, f"校验进程异常: exit {result.returncode}"

        output = result.stdout.decode("utf-8", errors="replace").strip()
        if not output:
            return False, "校验脚本无输出"

        resp = json.loads(output)
        if resp.get("valid"):
            return True, ""
        else:
            return False, resp.get("error", "未知错误")

    except subprocess.TimeoutExpired:
        return False, "校验超时 (>15s)"
    except json.JSONDecodeError:
        raw = result.stdout.decode("utf-8", errors="replace")[:200] if result.stdout else ""
        return False, f"校验输出解析失败: {raw}"
    except Exception as e:
        return False, f"校验异常: {e}"


# OpenAI SSE delta 标准字段：role/content/tool_calls/finish_reason/index/logprobs
# 其余 string 字段均视为 reasoning/thinking 内容（如 reasoning_content、reasoning、thinking 等）
_STANDARD_DELTA_FIELDS = frozenset({"role", "content", "tool_calls", "finish_reason", "index", "logprobs"})


def _collect_sse_stream(resp, thinking_budget: int = None, log_prefix: str = "[image_to_desc]", start_time: float = None):
    """从 SSE 流式响应中收集 content 和 reasoning（通用收集非标准 delta 字段）。

    resp.encoding 必须由调用方预先设置为 'utf-8'。

    返回 (content: str, reasoning_len: int, chunk_count: int, truncated: bool)。
    truncated=True 表示 reasoning 超过预算被本地截断，调用方应返回 None。
    """
    content_parts = []
    reasoning_parts = []
    reasoning_char_count = 0
    chunk_count = 0

    for line in resp.iter_lines(decode_unicode=True):
        if not line or not line.startswith("data: "):
            continue
        data_str = line[6:]  # 去掉 "data: " 前缀
        if data_str.strip() == "[DONE]":
            break
        try:
            chunk = json.loads(data_str)
        except json.JSONDecodeError:
            continue
        choices = chunk.get("choices", [])
        if not choices:
            continue
        delta = choices[0].get("delta", {})
        chunk_count += 1

        if chunk_count == 1 and start_time is not None:
            print(f"{log_prefix} 首个 chunk 到达: {time.time() - start_time:.1f}s", flush=True)

        # 收集思考过程（如有）— 通用策略：收集所有非标准 delta string 字段
        # 兼容 reasoning_content (OpenAI)、reasoning (Qwen3.6) 等不同字段名，无需维护字段名列表
        reasoning = "".join(
            v for k, v in delta.items()
            if k not in _STANDARD_DELTA_FIELDS and isinstance(v, str) and v
        )
        if reasoning:
            reasoning_parts.append(reasoning)
            reasoning_char_count += len(reasoning)
            # 本地截断：当 reasoning 字符数估算 token 超过预算时主动关闭连接
            if thinking_budget and thinking_budget > 0 and reasoning_char_count // 2 > thinking_budget:
                est_tokens = reasoning_char_count // 2
                print(f"{log_prefix} [thinking_limit] 本地截断: reasoning ~{est_tokens} tokens > budget {thinking_budget}", flush=True)
                resp.close()
                return "", 0, chunk_count, True

        # 收集正式回复
        content = delta.get("content")
        if content:
            content_parts.append(content)

    full_content = "".join(content_parts)
    reasoning_len = len("".join(reasoning_parts))
    return full_content, reasoning_len, chunk_count, False


def _call_vision_llm_raw(
    image_data_url: str,
    llm_url: str,
    llm_api_key: str,
    llm_model: str,
    system_prompt: str,
    user_prompt: str,
    timeout: int = 120,
    thinking_budget: int = None,
) -> Optional[str]:
    """调用 Vision LLM，返回原始文本内容（不做 JSON 解析）。

    使用流式输出（stream=True），逐步收集内容，避免长响应超时。
    当 thinking_budget > 0 时开启思考模式（enable_thinking=True），兼容非思考模型。
    thinking_budget 限制思考长度，避免简单任务过度推理。
    """
    _L = "[image_to_desc]"
    chat_url = _build_chat_url(llm_url)
    if not chat_url:
        print(f"{_L} _build_chat_url returned empty for url={llm_url}", flush=True)
        return None

    headers = {"Content-Type": "application/json"}
    if llm_api_key:
        headers["Authorization"] = f"Bearer {llm_api_key}"

    payload = {
        "model": llm_model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": user_prompt},
                    {"type": "image_url", "image_url": {"url": image_data_url}},
                ],
            },
        ],
        "temperature": 0.1,
        "max_tokens": 4096,
        "stream": True,
        "stream_options": {"include_usage": True},
    }

    # thinking_budget > 0 开启思考模式；=0 显式禁用（兼容 thinking 模型的非思考模式）
    if thinking_budget is not None and thinking_budget > 0:
        payload["enable_thinking"] = True
        payload["thinking_budget"] = thinking_budget
    else:
        payload["enable_thinking"] = False

    payload_size_kb = len(json.dumps(payload, ensure_ascii=False)) / 1024
    _budget_info = f", thinking_budget={thinking_budget}" if (thinking_budget is not None and thinking_budget > 0) else ""
    print(f"{_L} Vision API request (stream): url={chat_url}, model={llm_model}, "
          f"payload_size={payload_size_kb:.1f} KB, timeout={timeout}s{_budget_info}", flush=True)
    print(f"{_L}   system_prompt[:200]: {system_prompt[:200]}", flush=True)
    print(f"{_L}   user_prompt[:200]: {user_prompt[:200]}", flush=True)

    start_time = time.time()
    try:
        resp = requests.post(chat_url, headers=headers, json=payload, timeout=(10, timeout), stream=True)

        if not resp.ok:
            elapsed = time.time() - start_time
            resp_preview = resp.text[:500] if resp.text else "(empty)"
            print(f"{_L} Vision API failed: status={resp.status_code}, elapsed={elapsed:.1f}s, "
                  f"body={resp_preview}", flush=True)
            return None

        # 强制 UTF-8 解码，避免 LLM API 未声明 charset 时回退到 latin-1 导致乱码
        resp.encoding = 'utf-8'

        full_content, reasoning_len, chunk_count, truncated = _collect_sse_stream(
            resp, thinking_budget=thinking_budget, log_prefix=_L, start_time=start_time
        )
        if truncated:
            return None

        elapsed = time.time() - start_time
        print(f"{_L} Vision API stream done: elapsed={elapsed:.1f}s, "
              f"chunks={chunk_count}, content_len={len(full_content)}, "
              f"reasoning_len={reasoning_len}", flush=True)

        if not full_content:
            print(f"{_L} Vision API returned empty content", flush=True)
            return None

        print(f"{_L} LLM content preview: {full_content[:2000]}", flush=True)
        return full_content

    except requests.exceptions.ConnectTimeout:
        elapsed = time.time() - start_time
        print(f"{_L} Vision API connect timeout: elapsed={elapsed:.1f}s", flush=True)
        return None
    except requests.exceptions.ReadTimeout:
        elapsed = time.time() - start_time
        print(f"{_L} Vision API read timeout: elapsed={elapsed:.1f}s", flush=True)
        return None
    except requests.exceptions.ConnectionError as e:
        elapsed = time.time() - start_time
        print(f"{_L} Vision API connection error: {e}", flush=True)
        return None
    except Exception as e:
        elapsed = time.time() - start_time
        print(f"{_L} Vision API error: type={type(e).__name__}, error={e}", flush=True)
        return None


def _call_vision_llm(
    image_data_url: str,
    llm_url: str,
    llm_api_key: str,
    llm_model: str,
    timeout: int = 120,
) -> Optional[dict]:
    """三步策略编排器：分类 → Plan A/B → 兜底降级。

    流程:
    1. 调用分类接口判断图片复杂度和是否为图表
    2. 简单图表 → Plan A (JSON 结构化 → 代码生成)
    3. 复杂图表 → Plan B (直接生成 Mermaid)，最多重试 2 次
    4. Plan B 全部失败 → 兜底降级为 Plan A (简化 JSON)

    返回的 dict 包含:
    - is_diagram, diagram_type, description, complexity, route
    - mermaid_code: 直接可用的 Mermaid 代码（当 is_diagram 为 True 时）
    """
    _L = "[image_to_desc]"

    # ===== Step 1: 复杂度分类 =====
    print(f"{_L} Step 1: 分类复杂度...", flush=True)
    classify_raw = _call_vision_llm_raw(
        image_data_url, llm_url, llm_api_key, llm_model,
        CLASSIFY_SYSTEM_PROMPT, CLASSIFY_USER_PROMPT, timeout,
        thinking_budget=get_thinking_limit("vision_classify", 2048),
    )
    if classify_raw is None:
        print(f"{_L} 分类调用失败 (raw is None)", flush=True)
        return None

    print(f"{_L} 分类 raw[:1000]: {classify_raw[:1000]}", flush=True)

    classify_result = _parse_llm_json_response(classify_raw)
    if classify_result is None:
        print(f"{_L} 分类结果解析失败, raw[:500]={classify_raw[:500]}", flush=True)
        return None

    is_diagram = classify_result.get("is_diagram", False)
    complexity = classify_result.get("complexity", "simple")
    diagram_type = classify_result.get("diagram_type", "")
    description = classify_result.get("description", "")
    node_est = classify_result.get("node_count_estimate", 0)

    print(f"{_L} 分类结果: is_diagram={is_diagram}, complexity={complexity}, "
          f"type={diagram_type}, nodes≈{node_est}", flush=True)

    if not is_diagram:
        return {
            "is_diagram": False, "diagram_type": "", "description": description,
            "complexity": "", "route": "classify", "mermaid_code": None,
        }

    # ===== Step 2a: Plan A — 简单图表走 JSON 结构化 =====
    if complexity == "simple":
        print(f"{_L} Step 2a: Plan A (简单图表 → JSON)", flush=True)
        json_raw = _call_vision_llm_raw(
            image_data_url, llm_url, llm_api_key, llm_model,
            SYSTEM_PROMPT, USER_PROMPT, timeout,
            thinking_budget=get_thinking_limit("vision_mermaid", 4096),
        )
        if json_raw is not None:
            print(f"{_L} Plan A raw[:1000]: {json_raw[:1000]}", flush=True)
            json_result = _parse_llm_json_response(json_raw)
            if json_result is not None:
                print(f"{_L} Plan A JSON parsed: diagram_type={json_result.get('diagram_type')}, "
                      f"is_diagram={json_result.get('is_diagram')}, "
                      f"has_data={'diagram_data' in json_result}", flush=True)
                mermaid_code = _json_to_mermaid(json_result)
                if mermaid_code:
                    print(f"{_L} Plan A generated mermaid ({len(mermaid_code)} chars): {mermaid_code[:500]}", flush=True)
                    # Plan A 也需要校验生成的代码
                    valid, error = _validate_mermaid(mermaid_code)
                    if valid:
                        print(f"{_L} Plan A 成功: {len(mermaid_code)} 字符", flush=True)
                        return {
                            "is_diagram": True,
                            "diagram_type": json_result.get("diagram_type", diagram_type),
                            "description": json_result.get("description", description),
                            "complexity": "simple",
                            "route": "plan_a",
                            "mermaid_code": mermaid_code,
                        }
                    else:
                        print(f"{_L} Plan A 校验失败: error={error}, code[:300]={mermaid_code[:300]}", flush=True)
                else:
                    print(f"{_L} Plan A _json_to_mermaid 返回 None, json_result={json_result}", flush=True)
            else:
                print(f"{_L} Plan A JSON 解析失败, raw[:500]={json_raw[:500]}", flush=True)
        else:
            print(f"{_L} Plan A LLM 调用返回 None", flush=True)
        print(f"{_L} Plan A 失败，尝试 Plan B...", flush=True)
        # Plan A 失败时也尝试 Plan B

    # ===== Step 2b: Plan B — 复杂图表直接生成 Mermaid =====
    print(f"{_L} Step 2b: Plan B (直接生成 Mermaid)...", flush=True)

    # 首次生成
    direct_raw = _call_vision_llm_raw(
        image_data_url, llm_url, llm_api_key, llm_model,
        DIRECT_MERMAID_SYSTEM_PROMPT, DIRECT_MERMAID_USER_PROMPT, timeout,
        thinking_budget=get_thinking_limit("vision_mermaid", 4096),
    )
    if direct_raw:
        print(f"{_L} Plan B raw[:1000]: {direct_raw[:1000]}", flush=True)
    mermaid_code = _extract_mermaid_from_text(direct_raw) if direct_raw else None
    if mermaid_code:
        print(f"{_L} Plan B 提取到 mermaid ({len(mermaid_code)} chars): {mermaid_code[:500]}", flush=True)
        valid, error = _validate_mermaid(mermaid_code)
        if valid:
            print(f"{_L} Plan B 首次生成即成功: {len(mermaid_code)} 字符", flush=True)
            return {
                "is_diagram": True,
                "diagram_type": diagram_type,
                "description": description,
                "complexity": complexity,
                "route": "plan_b",
                "mermaid_code": mermaid_code,
            }
        else:
            print(f"{_L} Plan B 首次生成校验失败: error={error}, code[:300]={mermaid_code[:300]}", flush=True)
    else:
        error = "未提取到 mermaid 代码块"
        mermaid_code = ""
        print(f"{_L} Plan B 首次生成: {error}, raw[:500]={direct_raw[:500] if direct_raw else 'None'}", flush=True)

    # 纠错重试: 最多 2 次，把报错 + 错误代码发给 LLM 修正
    for fix_attempt in range(1, 3):
        print(f"{_L} Plan B 纠错重试 {fix_attempt}/2...", flush=True)
        fix_user_prompt = _build_fix_user_prompt(error, mermaid_code or "(无有效代码)")
        print(f"{_L}   fix_user_prompt[:300]: {fix_user_prompt[:300]}", flush=True)
        fix_raw = _call_vision_llm_raw(
            image_data_url, llm_url, llm_api_key, llm_model,
            FIX_MERMAID_SYSTEM_PROMPT, fix_user_prompt, timeout,
            thinking_budget=get_thinking_limit("vision_mermaid", 4096),
        )
        if fix_raw:
            print(f"{_L}   fix raw[:1000]: {fix_raw[:1000]}", flush=True)
        mermaid_code = _extract_mermaid_from_text(fix_raw) if fix_raw else None
        if mermaid_code:
            print(f"{_L}   fix 提取到 mermaid ({len(mermaid_code)} chars): {mermaid_code[:500]}", flush=True)
            valid, error = _validate_mermaid(mermaid_code)
            if valid:
                print(f"{_L} Plan B 纠错成功 (fix_attempt {fix_attempt}): {len(mermaid_code)} 字符", flush=True)
                return {
                    "is_diagram": True,
                    "diagram_type": diagram_type,
                    "description": description,
                    "complexity": complexity,
                    "route": "plan_b_fix",
                    "mermaid_code": mermaid_code,
                }
            else:
                print(f"{_L} Plan B 纠错 {fix_attempt} 仍失败: error={error}, code[:300]={mermaid_code[:300]}", flush=True)
        else:
            error = "未提取到 mermaid 代码块"
            print(f"{_L} Plan B 纠错 {fix_attempt}: {error}, raw[:500]={fix_raw[:500] if fix_raw else 'None'}", flush=True)

    print(f"{_L} Plan B 全部失败 (1次生成 + 2次纠错)", flush=True)

    # ===== Step 3: 兜底降级为 Plan A =====
    print(f"{_L} Step 3: 兜底降级 — 强制 Plan A (简化 JSON)...", flush=True)
    fallback_raw = _call_vision_llm_raw(
        image_data_url, llm_url, llm_api_key, llm_model,
        SYSTEM_PROMPT, FALLBACK_JSON_USER_PROMPT, timeout,
        thinking_budget=get_thinking_limit("vision_mermaid", 4096),
    )
    if fallback_raw is not None:
        print(f"{_L} 兜底 raw[:1000]: {fallback_raw[:1000]}", flush=True)
        fallback_result = _parse_llm_json_response(fallback_raw)
        if fallback_result is not None:
            print(f"{_L} 兜底 JSON parsed: {fallback_result}", flush=True)
            mermaid_code = _json_to_mermaid(fallback_result)
            if mermaid_code:
                print(f"{_L} 兜底 Plan A 成功: {len(mermaid_code)} 字符", flush=True)
                return {
                    "is_diagram": True,
                    "diagram_type": fallback_result.get("diagram_type", diagram_type),
                    "description": fallback_result.get("description", description),
                    "complexity": complexity,
                    "route": "fallback",
                    "mermaid_code": mermaid_code,
                }
            else:
                print(f"{_L} 兜底 _json_to_mermaid 返回 None, json_result={fallback_result}", flush=True)
        else:
            print(f"{_L} 兜底 JSON 解析失败, raw[:500]={fallback_raw[:500]}", flush=True)
    else:
        print(f"{_L} 兜底 LLM 调用返回 None", flush=True)

    print(f"{_L} 所有方案均失败", flush=True)
    return None


def _extract_mermaid_code(llm_result: dict) -> Optional[str]:
    """从 LLM 结果 dict 中提取 Mermaid 代码。

    优先使用 result 中已有的 mermaid_code 字段（三步策略直接生成），
    否则回退到从 diagram_data JSON 生成。
    """
    if not llm_result:
        return None

    # 三步策略已在 dict 中放入了 mermaid_code
    mc = llm_result.get("mermaid_code")
    if mc:
        return mc

    # 兼容旧逻辑: 从结构化 JSON 生成
    return _json_to_mermaid(llm_result)


def scan_images(assets_dir: str) -> List[str]:
    """Scan assets directory for image files."""
    images = []
    if not os.path.isdir(assets_dir):
        return images

    for root, dirs, files in os.walk(assets_dir):
        for fname in sorted(files):
            ext = os.path.splitext(fname)[1].lower()
            if ext in IMAGE_EXTENSIONS:
                images.append(os.path.join(root, fname))

    return images


def process_image(
    image_path: str,
    llm_url: str,
    llm_api_key: str,
    llm_model: str,
) -> Dict[str, Any]:
    """Process a single image: identify and generate Mermaid code."""
    result = {
        "image": image_path,
        "filename": os.path.basename(image_path),
        "status": "skipped",
        "diagram_type": "",
        "description": "",
        "mermaid_file": "",
        "error": "",
    }

    # Check if mermaid file already exists
    mermaid_path = image_path + ".mermaid"
    meta_path = image_path + ".mermaid.json"
    if os.path.exists(mermaid_path):
        result["status"] = "already_processed"
        result["mermaid_file"] = mermaid_path
        return result

    # Encode image
    data_url = _encode_image_to_data_url(image_path)
    if not data_url:
        result["status"] = "error"
        result["error"] = "Failed to encode image"
        return result

    # Call Vision LLM
    llm_result = _call_vision_llm(data_url, llm_url, llm_api_key, llm_model)
    if llm_result is None:
        result["status"] = "error"
        result["error"] = "Vision API call failed or model not supported"
        return result

    result["description"] = llm_result.get("description", "")
    result["diagram_type"] = llm_result.get("diagram_type", "")

    # Extract Mermaid code
    mermaid_code = _extract_mermaid_code(llm_result)
    if not mermaid_code:
        result["status"] = "not_diagram"
        # Save metadata anyway for reference
        meta = {
            "image": os.path.basename(image_path),
            "is_diagram": False,
            "description": result["description"],
            "processed_at": datetime.now().isoformat(),
        }
        try:
            with open(meta_path, "w", encoding="utf-8") as f:
                json.dump(meta, f, ensure_ascii=False, indent=2)
        except Exception:
            pass
        return result

    # Save Mermaid code
    try:
        with open(mermaid_path, "w", encoding="utf-8") as f:
            f.write(mermaid_code)

        # Save metadata
        meta = {
            "image": os.path.basename(image_path),
            "is_diagram": True,
            "diagram_type": result["diagram_type"],
            "description": result["description"],
            "mermaid_file": os.path.basename(mermaid_path),
            "processed_at": datetime.now().isoformat(),
        }
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump(meta, f, ensure_ascii=False, indent=2)

        result["status"] = "success"
        result["mermaid_file"] = mermaid_path
        print(f"[image_to_desc] Generated Mermaid for {os.path.basename(image_path)}: {result['diagram_type']}")

    except Exception as e:
        result["status"] = "error"
        result["error"] = f"Failed to save: {e}"

    return result


_INJECT_MARKER_PREFIX = "<!-- mermaid-injected:"
_INJECT_MARKER_SUFFIX = "-->"

# Regex to match markdown image references pointing into raw/assets/
_IMAGE_REF_RE = re.compile(
    r'!\[([^\]]*)\]\([^)]*(?:raw/assets/|/raw/assets/)([^)]+)\)'
)


def _build_inject_marker(image_name: str) -> str:
    return f"{_INJECT_MARKER_PREFIX} {image_name} {_INJECT_MARKER_SUFFIX}"


def inject_mermaid_into_wiki(project_path: str) -> Dict[str, Any]:
    """Scan wiki/ pages and inject Mermaid code blocks after matching image refs.

    For every ``![...](../../raw/assets/foo.png)`` whose ``foo.png.mermaid``
    exists in ``raw/assets/``, a `````mermaid````` block is inserted right after
    the image reference.  A marker comment is added so repeated runs are
    idempotent and the injection can be cleaned up later.

    Failure handling
    ----------------
    * Wiki dir missing  → return early with ``injected=0``.
    * Cannot read/write a wiki page → log warning, skip that file.
    * Mermaid file empty → skip that image.
    * Image already has an injected block (marker present) → skip.
    * No mermaid files found → return early with ``injected=0``.
    """
    wiki_dir = os.path.join(project_path, "wiki")
    assets_dir = os.path.join(project_path, "raw", "assets")

    result = {
        "ok": True,
        "scanned_pages": 0,
        "injected": 0,
        "skipped_already": 0,
        "errors": 0,
        "details": [],
    }

    if not os.path.isdir(wiki_dir):
        result["details"].append("wiki/ directory not found, nothing to inject")
        return result

    # ------------------------------------------------------------------
    # 1. Load mermaid map: {image_filename: mermaid_code}
    # ------------------------------------------------------------------
    mermaid_map: Dict[str, str] = {}
    if os.path.isdir(assets_dir):
        for fname in os.listdir(assets_dir):
            if not fname.endswith(".mermaid"):
                continue
            image_name = fname[: -len(".mermaid")]
            mermaid_path = os.path.join(assets_dir, fname)
            try:
                with open(mermaid_path, "r", encoding="utf-8") as fh:
                    code = fh.read().strip()
                if code:
                    mermaid_map[image_name] = code
                else:
                    print(f"[inject_mermaid] Empty mermaid file, skip: {fname}")
            except Exception as exc:
                print(f"[inject_mermaid] Failed to read {mermaid_path}: {exc}")

    if not mermaid_map:
        result["details"].append("No valid mermaid files found in raw/assets/")
        return result

    # ------------------------------------------------------------------
    # 2. Walk wiki/ pages and inject
    # ------------------------------------------------------------------
    for root, _dirs, files in os.walk(wiki_dir):
        for fname in sorted(files):
            if not fname.endswith(".md"):
                continue

            page_path = os.path.join(root, fname)
            rel_page = os.path.relpath(page_path, project_path).replace("\\", "/")
            result["scanned_pages"] += 1

            try:
                with open(page_path, "r", encoding="utf-8") as fh:
                    content = fh.read()
            except Exception as exc:
                print(f"[inject_mermaid] Failed to read {page_path}: {exc}")
                result["errors"] += 1
                result["details"].append(f"read_error: {rel_page}: {exc}")
                continue

            original = content
            page_injected = 0
            page_skipped = 0

            # Collect all matches first, then process in reverse order
            # so that insertions don't invalidate earlier match positions.
            matches = list(_IMAGE_REF_RE.finditer(content))

            for match in reversed(matches):
                image_name = match.group(2).strip()
                if image_name not in mermaid_map:
                    continue

                marker = _build_inject_marker(image_name)

                # Check if already injected (look in a window after the match)
                window_start = match.end()
                window_end = min(window_start + 500, len(content))
                window_text = content[window_start:window_end]
                if marker in window_text:
                    page_skipped += 1
                    continue

                # Also skip if there is already a ```mermaid block right after
                after = content[window_start:window_end].lstrip("\n")
                if after.startswith("```mermaid"):
                    page_skipped += 1
                    continue

                # Build the injection text
                mermaid_code = mermaid_map[image_name]
                injection = (
                    f"\n\n{marker}\n"
                    f"```mermaid\n{mermaid_code}\n```\n"
                )

                # Insert right after the image reference.
                content = (
                    content[: match.end()] + injection + content[match.end() :]
                )
                page_injected += 1

            if content == original:
                result["skipped_already"] += page_skipped
                continue

            try:
                with open(page_path, "w", encoding="utf-8") as fh:
                    fh.write(content)
            except Exception as exc:
                print(f"[inject_mermaid] Failed to write {page_path}: {exc}")
                result["errors"] += 1
                result["details"].append(f"write_error: {rel_page}: {exc}")
                continue

            result["injected"] += page_injected
            result["skipped_already"] += page_skipped
            if page_injected > 0:
                print(
                    f"[inject_mermaid] {rel_page}: "
                    f"injected {page_injected}, skipped {page_skipped}"
                )

    print(
        f"[inject_mermaid] Done: scanned={result['scanned_pages']}, "
        f"injected={result['injected']}, skipped={result['skipped_already']}, "
        f"errors={result['errors']}"
    )
    return result


def process_project_images(
    project_path: str,
    llm_url: str,
    llm_api_key: str,
    llm_model: str,
    force: bool = False,
) -> Dict[str, Any]:
    """Process all images in a project's raw/assets directory."""
    assets_dir = os.path.join(project_path, "raw", "assets")
    images = scan_images(assets_dir)

    if not images:
        return {
            "ok": True,
            "total": 0,
            "processed": 0,
            "success": 0,
            "not_diagram": 0,
            "errors": 0,
            "results": [],
            "message": "No images found in raw/assets",
        }

    results = []
    success_count = 0
    not_diagram_count = 0
    error_count = 0
    skipped_count = 0

    # 收集需要处理的图片（跳过已处理的）
    pending_images = []
    for image_path in images:
        if force:
            mermaid_path = image_path + ".mermaid"
            meta_path = image_path + ".mermaid.json"
            if os.path.exists(mermaid_path):
                os.remove(mermaid_path)
            if os.path.exists(meta_path):
                os.remove(meta_path)

        # 检查是否已处理
        mermaid_path = image_path + ".mermaid"
        if os.path.exists(mermaid_path):
            result = {
                "image": image_path,
                "filename": os.path.basename(image_path),
                "status": "already_processed",
                "diagram_type": "",
                "description": "",
                "mermaid_file": mermaid_path,
                "error": "",
            }
            results.append(result)
            skipped_count += 1
            continue

        pending_images.append(image_path)

    # 并发调用 Vision LLM
    if pending_images:
        print(f"[image_to_desc] Processing {len(pending_images)} image(s) concurrently (max_workers=4)...")

        def _process_one(img_path):
            return img_path, process_image(img_path, llm_url, llm_api_key, llm_model)

        with ThreadPoolExecutor(max_workers=4) as executor:
            futures = {executor.submit(_process_one, img): img for img in pending_images}
            for future in as_completed(futures):
                img_path = futures[future]
                try:
                    _, result = future.result()
                    results.append(result)
                    status = result["status"]
                    print(f"[image_to_desc] Done: {os.path.basename(img_path)} → {status}")
                except Exception as e:
                    print(f"[image_to_desc] EXCEPTION: {os.path.basename(img_path)}: {e}")
                    results.append({
                        "image": img_path,
                        "filename": os.path.basename(img_path),
                        "status": "error",
                        "diagram_type": "",
                        "description": "",
                        "mermaid_file": "",
                        "error": str(e),
                    })

    # 统计结果
    for result in results:
        if result["status"] == "success":
            success_count += 1
        elif result["status"] == "not_diagram":
            not_diagram_count += 1
        elif result["status"] == "error":
            error_count += 1
        elif result["status"] == "already_processed":
            skipped_count += 1

    # After all mermaid files are generated, inject them into wiki pages
    inject_result = {}
    try:
        inject_result = inject_mermaid_into_wiki(project_path)
    except Exception as exc:
        print(f"[image_to_desc] inject_mermaid_into_wiki failed: {exc}")
        inject_result = {"ok": False, "error": str(exc)}

    return {
        "ok": True,
        "total": len(images),
        "processed": success_count + not_diagram_count + error_count,
        "success": success_count,
        "not_diagram": not_diagram_count,
        "errors": error_count,
        "skipped": skipped_count,
        "results": results,
        "inject": inject_result,
    }


# ============================================================================
#  Markdown image → Mermaid (inline insertion)
# ============================================================================

# Regex: ![alt text](path) or ![alt text](path "title")
_MD_IMAGE_RE = re.compile(r'!\[([^\]]*)\]\(([^)\s]+)(?:\s+"[^"]*")?\)')

_MERMAID_INLINE_MARKER = "<!-- mermaid-inserted -->"


def _resolve_image_path(image_ref: str, md_dir: str, project_dir: str = "") -> Optional[str]:
    """Resolve a markdown image reference to an absolute path.

    Tries in order:
      1. If image_ref is already absolute → use as-is
      2. Relative to the markdown file's directory
      3. Relative to the project root (if provided)
      4. Relative to project root + raw/assets/ (KMA project structure)

    Returns the resolved path if file exists, else None.
    """
    # Strip URL fragment
    image_ref = image_ref.split("#")[0]

    # 1. Absolute path
    if os.path.isabs(image_ref):
        if os.path.isfile(image_ref):
            return image_ref
        return None

    # Normalize forward slashes
    norm_ref = os.path.normpath(image_ref)

    # 2. Relative to markdown file directory
    candidate = os.path.normpath(os.path.join(md_dir, norm_ref))
    if os.path.isfile(candidate):
        return candidate

    # 3. Relative to project root
    if project_dir:
        candidate = os.path.normpath(os.path.join(project_dir, norm_ref))
        if os.path.isfile(candidate):
            return candidate

        # 4. project_root/raw/assets/
        candidate = os.path.normpath(os.path.join(project_dir, "raw", "assets", norm_ref))
        if os.path.isfile(candidate):
            return candidate

    return None


def _md_log(msg: str):
    """Log with flush so output is visible in real-time from subprocess."""
    print(msg, flush=True)


def process_markdown_images(
    md_file_path: str,
    llm_url: str,
    llm_api_key: str,
    llm_model: str,
) -> Dict[str, Any]:
    """Parse a markdown file, detect diagram images, and insert mermaid code in-place.

    Unlike ``process_project_images`` which scans raw/assets/ and writes separate
    .mermaid files, this function:
      1. Finds image references inside the markdown itself
      2. Resolves each image path (relative to the md file, or absolute)
      3. Calls Vision LLM to check if it's a diagram
      4. If yes, inserts a `````mermaid`` block right after the image reference
      5. Writes the modified markdown back to disk
    """
    _L = "[md_mermaid]"
    log = lambda msg: _md_log(f"{_L} {msg}")
    result: Dict[str, Any] = {
        "ok": True,
        "md_file": md_file_path,
        "total_images": 0,
        "diagrams_found": 0,
        "inserted": 0,
        "errors": 0,
        "skipped_urls": 0,
        "skipped_not_found": 0,
        "details": [],
    }

    log(f"====== START: {md_file_path}")

    if not os.path.isfile(md_file_path):
        log(f"ERROR: markdown file not found: {md_file_path}")
        result["ok"] = False
        result["error"] = f"Markdown file not found: {md_file_path}"
        return result

    md_abs_path = os.path.abspath(md_file_path)
    md_dir = os.path.dirname(md_abs_path)
    # Infer project root: go up from md_dir to find the project directory
    # (md files are typically under project/sources/ or project/sources/subfolder/)
    project_dir = ""
    probe = md_dir
    for _ in range(5):
        if os.path.isdir(os.path.join(probe, "raw", "assets")):
            project_dir = probe
            break
        parent = os.path.dirname(probe)
        if parent == probe:
            break
        probe = parent

    log(f"md_file (abs)   : {md_abs_path}")
    log(f"md_dir          : {md_dir}")
    log(f"project_dir     : {project_dir or '(not detected)'}")

    try:
        with open(md_abs_path, "r", encoding="utf-8") as fh:
            content = fh.read()
    except Exception as exc:
        log(f"ERROR reading file: {exc}")
        result["ok"] = False
        result["error"] = f"Failed to read markdown: {exc}"
        return result

    log(f"file size       : {len(content)} chars")

    original = content
    matches = list(_MD_IMAGE_RE.finditer(content))

    if not matches:
        log(f"No image references found in markdown")
        result["details"].append("No image references found")
        return result

    result["total_images"] = len(matches)
    log(f"Found {len(matches)} image reference(s):")
    for i, m in enumerate(matches):
        log(f"  [{i+1}] alt='{m.group(1).strip()[:40]}' ref='{m.group(2).strip()[:80]}'")

    # ---- Phase 1: 串行预处理 — 解析路径、编码图片、过滤跳过项 ----
    # 收集需要调用 LLM 的任务
    pending_tasks = []  # [{match, image_ref, img_name, data_url, img_size_kb}]
    for idx, match in enumerate(matches):
        image_ref = match.group(2).strip()
        alt_text = match.group(1).strip()
        img_name = os.path.basename(image_ref.split("#")[0].split("?")[0])

        # Skip remote URLs and inline data URIs
        if image_ref.startswith(("http://", "https://", "data:")):
            log(f"  [{idx+1}] SKIP: remote URL or data URI ({image_ref[:60]})")
            result["skipped_urls"] += 1
            continue

        # Resolve image path
        image_path = _resolve_image_path(image_ref, md_dir, project_dir)
        if image_path is None:
            log(f"  [{idx+1}] SKIP: image not found ({image_ref[:60]})")
            result["skipped_not_found"] += 1
            result["details"].append(f"Image not found: {image_ref}")
            continue

        # Check if mermaid already inserted
        check_window = content[match.end():match.end() + len(_MERMAID_INLINE_MARKER) + 2]
        if _MERMAID_INLINE_MARKER in check_window:
            log(f"  [{idx+1}] SKIP: already has mermaid marker ({img_name})")
            result["details"].append(f"Already inserted: {img_name}")
            continue

        # Encode image
        data_url = _encode_image_to_data_url(image_path)
        if not data_url:
            log(f"  [{idx+1}] ERROR: failed to encode ({img_name})")
            result["errors"] += 1
            result["details"].append(f"Encode failed: {image_ref}")
            continue

        img_size_kb = os.path.getsize(image_path) / 1024
        log(f"  [{idx+1}] resolved: {image_path} ({img_size_kb:.1f} KB)")
        pending_tasks.append({
            "match": match,
            "image_ref": image_ref,
            "img_name": img_name,
            "data_url": data_url,
            "img_size_kb": img_size_kb,
        })

    # ---- Phase 2: 并发调用 Vision LLM ----
    llm_results = {}  # match_id → (llm_result, elapsed)
    if pending_tasks:
        log(f"Calling Vision LLM for {len(pending_tasks)} image(s) concurrently (max_workers=4)...")

        def _call_llm_task(task):
            match_id = id(task["match"])
            t0 = time.time()
            r = _call_vision_llm(task["data_url"], llm_url, llm_api_key, llm_model)
            return match_id, r, time.time() - t0

        with ThreadPoolExecutor(max_workers=4) as executor:
            futures = {executor.submit(_call_llm_task, t): t for t in pending_tasks}
            for future in as_completed(futures):
                task = futures[future]
                try:
                    match_id, llm_result, elapsed = future.result()
                    llm_results[id(task["match"])] = (llm_result, elapsed)
                    if llm_result is None:
                        log(f"  → ERROR: Vision API returned None ({elapsed:.1f}s) for {task['img_name']}")
                    else:
                        is_diagram = llm_result.get("is_diagram", False)
                        diagram_type = llm_result.get("diagram_type", "")
                        desc = llm_result.get("description", "")[:60]
                        log(f"  → LLM result ({elapsed:.1f}s): {task['img_name']}: is_diagram={is_diagram}, type={diagram_type}, desc='{desc}'")
                except Exception as e:
                    log(f"  → EXCEPTION: {task['img_name']}: {e}")
                    llm_results[id(task["match"])] = (None, 0)

    # ---- Phase 3: 逆序插入 mermaid 代码（避免偏移） ----
    for task in reversed(pending_tasks):
        match = task["match"]
        img_name = task["img_name"]
        image_ref = task["image_ref"]

        llm_result, elapsed = llm_results.get(id(match), (None, 0))
        if llm_result is None:
            result["errors"] += 1
            result["details"].append(f"Vision API failed: {image_ref}")
            continue

        is_diagram = llm_result.get("is_diagram", False)
        diagram_type = llm_result.get("diagram_type", "")
        description = llm_result.get("description", "")[:80]

        mermaid_code = _extract_mermaid_code(llm_result)
        if not mermaid_code:
            log(f"  → SKIP: not a diagram or no mermaid code ({img_name})")
            result["details"].append(f"Not a diagram: {img_name} ({description})")
            continue

        mermaid_lines = mermaid_code.count("\n") + 1
        log(f"  → MERMAID OK: {img_name}: type={diagram_type}, {mermaid_lines} lines, inserting...")

        injection = (
            f"\n\n{_MERMAID_INLINE_MARKER}\n"
            f"```mermaid\n{mermaid_code}\n```\n"
        )
        content = content[:match.end()] + injection + content[match.end():]
        result["inserted"] += 1
        result["diagrams_found"] += 1
        result["details"].append(f"Inserted {diagram_type} mermaid: {img_name}")

    # Write back only if content changed
    if content != original:
        try:
            with open(md_abs_path, "w", encoding="utf-8") as fh:
                fh.write(content)
            log(f"Written updated file: {md_abs_path}")
        except Exception as exc:
            log(f"ERROR writing file: {exc}")
            result["ok"] = False
            result["error"] = f"Failed to write back: {exc}"
            return result
    else:
        log(f"No changes to write back")

    log(
        f"====== DONE: total={result['total_images']}, "
        f"diagrams={result['diagrams_found']}, inserted={result['inserted']}, "
        f"errors={result['errors']}, skipped_urls={result['skipped_urls']}, "
        f"skipped_not_found={result['skipped_not_found']}"
    )
    return result


# ============================================================================
#  New Pipeline: File → Structured Content (Mermaid, Table, Code, etc.)
#  Uses prefix cache strategy + streaming + thinking
# ============================================================================

# Prompt: 图片描述 (配合文档上下文, 用于 prefix cache)
# 文档全文放在 system message 中，确保所有图片请求的 system 前缀完全相同，触发 API prefix cache
DESCRIBE_IMAGE_SYSTEM_PROMPT_TEMPLATE = """你是一个专业的图片内容描述专家。

## 任务
你将收到一份文档的完整文本（其中 {{IMAGE_N}} 是图片占位符），以及其中一张图片。
请用 200 字以内简洁描述这张图片的内容，重点关注图片中的结构化信息（如图表结构、表格数据、代码逻辑等）。

## 输出规则
- 直接输出描述文本，不要添加任何前缀、标记或额外文字
- 200 字以内
- 如果图片包含图表，请说明图表类型和主要内容
- 如果图片包含表格，请说明表格结构和列名
- 如果图片包含代码，请说明编程语言和代码功能
- 如果是照片或插图，简要描述内容

## 文档全文（图片用 {{IMAGE_N}} 占位符标记）
---
{context_text}
---"""

# user prompt 只包含简短指令 + 图片索引，不含文档全文（全文在 system 中，用于 prefix cache）
DESCRIBE_IMAGE_USER_PROMPT_TEMPLATE = "请描述标记为 {{IMAGE_{idx}}} 的这张图片的内容（200字以内）。"

# Prompt: 判断图片是否可转为结构化文本 + 生成代码
CONVERT_JUDGE_SYSTEM_PROMPT = """你是一个专业的图片内容结构化转换专家。

## 任务
根据图片内容和已有的描述，判断这张图片是否可以转换为结构化的文本格式。如果可以，生成对应的代码。

## 支持的转换类型

### mermaid（图表）
流程图、时序图、类图、状态图、ER 图、思维导图、旅程图、时间线图等。
输出 ```mermaid ... ``` 代码块。

### table（表格）
图片中包含表格数据、对比矩阵、参数列表等。
输出 Markdown 表格格式。

### code（代码）
图片中包含可辨识的源代码。
输出对应语言的代码块（如 ```python ... ```）。

### list（结构化列表）
图片中有清晰的层级列表、步骤、分类信息，适合用结构化列表表达。
输出 Markdown 列表。

## 不可转换的情况（返回 convertible: false）
- 纯装饰性图片、Logo、图标
- 照片（人物、风景等）
- 模糊不清的图片
- 纯文字截图（直接 OCR 即可，不需要结构化转换）
- 含合并单元格、跨行、跨列的表格（Markdown 表格无法表达）

## 输出格式
严格按以下 JSON 格式返回：

可转换时：
```json
{
  "convertible": true,
  "type": "mermaid/table/code/list",
  "description": "简要说明",
  "code": "生成的完整代码"
}
```

不可转换时：
```json
{
  "convertible": false,
  "type": "",
  "description": "简要说明图片内容",
  "code": ""
}
```

## Mermaid 代码规则
1. 节点 ID 只使用英文字母、数字和下划线
2. 禁止使用 Mermaid 保留关键字作为节点 ID：start, end, subgraph, if, else, loop, alt, opt, par
3. 包含特殊字符的标签用双引号包裹
4. 中文内容保留原样
5. 包含图中所有可见的节点和连接关系

## 重要
- 只输出一个 JSON 对象，不要添加任何额外文字"""

CONVERT_JUDGE_USER_PROMPT_TEMPLATE = """图片描述: {description}

请判断这张图片是否可以转换为结构化文本，如果可以，生成对应的代码。"""

# Prompt: 修复表格语法
FIX_TABLE_SYSTEM_PROMPT = """你是一个 Markdown 表格修复专家。上一次生成的表格格式有误，请修正。
只输出修正后的 Markdown 表格，不要有任何其他文字。"""


def _extract_text_and_placeholders(md_content: str) -> tuple:
    """从 Markdown 内容中提取文本和图片占位符。

    将图片引用替换为占位符 {{IMAGE_N}}，同时记录每个图片的信息。

    Returns:
        (text_with_placeholders, image_infos)
        - text_with_placeholders: 替换图片后的完整 markdown 文本
        - image_infos: [{
            "index": 1,
            "placeholder": "{{IMAGE_1}}",
            "alt_text": "...",
            "image_ref": "path/to/image.png",
            "line_number": 5,
          }, ...]
    """
    image_infos = []
    lines = md_content.split('\n')
    result_lines = []
    img_counter = 0

    for line_num, line in enumerate(lines):
        new_line = line
        # 查找本行所有图片引用
        for match in _MD_IMAGE_RE.finditer(line):
            img_counter += 1
            alt_text = match.group(1).strip()
            image_ref = match.group(2).strip()
            placeholder = f"{{{{IMAGE_{img_counter}}}}}"

            image_infos.append({
                "index": img_counter,
                "placeholder": placeholder,
                "alt_text": alt_text,
                "image_ref": image_ref,
                "line_number": line_num + 1,
            })

            new_line = new_line.replace(match.group(0), placeholder, 1)

        result_lines.append(new_line)

    text_with_placeholders = '\n'.join(result_lines)
    return text_with_placeholders, image_infos


def _call_vision_llm_with_context(
    image_data_url: str,
    context_text: str,
    user_prompt: str,
    llm_url: str,
    llm_api_key: str,
    llm_model: str,
    timeout: int = 120,
) -> Optional[str]:
    """带文档上下文的 Vision LLM 调用 (流式 + 思考)。

    messages 结构: [system(指令+全文), user(短文本+image)]
    文档全文放在 system message 中，所有图片请求的 system 前缀完全相同，
    可触发 API 侧 prefix cache（命中 system + 全文，占总 token 95%+）。
    """
    chat_url = _build_chat_url(llm_url)
    if not chat_url:
        return None

    headers = {"Content-Type": "application/json"}
    if llm_api_key:
        headers["Authorization"] = f"Bearer {llm_api_key}"

    # system message 包含指令 + 文档全文（所有请求完全相同，prefix cache 命中）
    system_content = DESCRIBE_IMAGE_SYSTEM_PROMPT_TEMPLATE.format(context_text=context_text)

    _desc_budget = get_thinking_limit("vision_describe", 4096)

    payload = {
        "model": llm_model,
        "messages": [
            {"role": "system", "content": system_content},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": user_prompt},
                    {"type": "image_url", "image_url": {"url": image_data_url}},
                ],
            },
        ],
        "temperature": 0.1,
        "max_tokens": 2048,
        "stream": True,
        "stream_options": {"include_usage": True},
    }
    if _desc_budget and _desc_budget > 0:
        payload["enable_thinking"] = True
        payload["thinking_budget"] = _desc_budget
    else:
        payload["enable_thinking"] = False

    start_time = time.time()
    try:
        resp = requests.post(
            chat_url, headers=headers, json=payload,
            timeout=(10, timeout), stream=True
        )
        if not resp.ok:
            print(f"[pipeline] Describe API failed: status={resp.status_code}", flush=True)
            return None

        # 强制 UTF-8 解码，避免 LLM API 未声明 charset 时回退到 latin-1 导致乱码
        resp.encoding = 'utf-8'

        full_content, reasoning_len, _chunk_count, truncated = _collect_sse_stream(
            resp, thinking_budget=_desc_budget, log_prefix="[pipeline]", start_time=start_time
        )
        if truncated:
            return None

        elapsed = time.time() - start_time
        print(f"[pipeline] Describe done: {elapsed:.1f}s, content_len={len(full_content)}, reasoning_len={reasoning_len}", flush=True)
        return full_content.strip() if full_content else None

    except Exception as e:
        print(f"[pipeline] Describe error: {e}", flush=True)
        return None


_UNCONVERTIBLE_TABLE_KEYWORDS = [
    "合并单元格", "跨行", "跨列",
    "merged cell", "rowspan", "colspan",
    "嵌套表格", "不规则表格",
]

def _is_unconvertible_table(description: str) -> bool:
    """检查描述是否表明表格含合并单元格等不可转换特征。"""
    if not description:
        return False
    desc_lower = description.lower()
    return any(kw in desc_lower for kw in _UNCONVERTIBLE_TABLE_KEYWORDS)


def _validate_table(code: str) -> tuple:
    """校验 Markdown 表格语法。

    返回 (is_valid: bool, error_msg: str)。
    """
    if not code or len(code.strip()) < 5:
        return False, "表格内容太短"

    lines = [l.strip() for l in code.strip().split('\n') if l.strip()]
    if len(lines) < 2:
        return False, "表格至少需要表头和数据行"

    # 检查表头
    if '|' not in lines[0]:
        return False, "表头缺少列分隔符 |"

    # 检查分隔行
    if len(lines) > 1:
        sep_line = lines[1]
        if not re.match(r'^[\|\s:\-]+$', sep_line):
            return False, "第二行应为分隔行 (如 |---|---|)"

    # 检查列数一致性
    header_cols = len([c for c in lines[0].split('|') if c.strip()])
    for i, line in enumerate(lines[2:], start=3):
        cols = len([c for c in line.split('|') if c.strip()])
        if cols != header_cols:
            return False, f"第{i}行列数({cols})与表头({header_cols})不一致"

    return True, ""


def _judge_and_convert_image(
    image_data_url: str,
    description: str,
    llm_url: str,
    llm_api_key: str,
    llm_model: str,
    timeout: int = 120,
) -> dict:
    """判断图片是否可转为结构化文本，并生成代码。

    对 mermaid 使用 CLI 校验 + 纠错重试（最多 2 次）。
    对 table 使用基础校验 + 纠错重试（最多 2 次）。

    Returns dict: {convertible, type, description, code, route}
    """
    _L = "[pipeline:convert]"

    user_prompt = CONVERT_JUDGE_USER_PROMPT_TEMPLATE.format(description=description)

    # Step 1: 初始判断 + 生成 (限制思考预算，分类任务不需要深度推理)
    raw = _call_vision_llm_raw(
        image_data_url, llm_url, llm_api_key, llm_model,
        CONVERT_JUDGE_SYSTEM_PROMPT, user_prompt, timeout,
        thinking_budget=get_thinking_limit("vision_convert", 4096),
    )
    if raw is None:
        return {"convertible": False, "type": "", "description": description,
                "code": "", "route": "error"}

    result = _parse_llm_json_response(raw)
    if result is None:
        return {"convertible": False, "type": "", "description": description,
                "code": "", "route": "parse_error"}

    content_type = result.get("type", "")
    convertible = result.get("convertible", False)
    code = result.get("code", "")
    desc = result.get("description", description)

    if not convertible or not code:
        return {"convertible": False, "type": "", "description": desc,
                "code": "", "route": "not_convertible"}

    # Step 2: 根据类型校验
    if content_type == "mermaid":
        # 从文本中提取 mermaid 代码块（如果 LLM 返回了带 ``` 的格式）
        extracted = _extract_mermaid_from_text(code)
        if extracted:
            code = extracted

        valid, error = _validate_mermaid(code)
        if valid:
            return {"convertible": True, "type": "mermaid", "description": desc,
                    "code": code, "route": "mermaid_ok"}

        # 纠错重试 (最多 2 次)
        for attempt in range(1, 3):
            print(f"{_L} Mermaid 纠错 {attempt}/2: {error}", flush=True)
            fix_prompt = _build_fix_user_prompt(error, code)
            fix_raw = _call_vision_llm_raw(
                image_data_url, llm_url, llm_api_key, llm_model,
                FIX_MERMAID_SYSTEM_PROMPT, fix_prompt, timeout,
                thinking_budget=get_thinking_limit("vision_convert", 4096),
            )
            fix_code = _extract_mermaid_from_text(fix_raw) if fix_raw else None
            if fix_code:
                valid, error = _validate_mermaid(fix_code)
                if valid:
                    return {"convertible": True, "type": "mermaid", "description": desc,
                            "code": fix_code, "route": f"mermaid_fix_{attempt}"}

        return {"convertible": False, "type": "mermaid", "description": desc,
                "code": "", "route": "mermaid_failed"}

    elif content_type == "table":
        valid, error = _validate_table(code)
        if valid:
            return {"convertible": True, "type": "table", "description": desc,
                    "code": code, "route": "table_ok"}

        # 纠错重试 (最多 2 次)
        for attempt in range(1, 3):
            print(f"{_L} Table 纠错 {attempt}/2: {error}", flush=True)
            fix_user = f"表格格式有误: {error}\n\n原始表格:\n{code}\n\n请修正表格格式。"
            fix_raw = _call_vision_llm_raw(
                image_data_url, llm_url, llm_api_key, llm_model,
                FIX_TABLE_SYSTEM_PROMPT, fix_user, timeout,
                thinking_budget=get_thinking_limit("vision_convert", 4096),
            )
            if fix_raw:
                fix_code = fix_raw.strip()
                # 尝试从代码块中提取
                block_match = re.search(r'```(?:markdown)?\s*\n?([\s\S]*?)```', fix_code)
                if block_match:
                    fix_code = block_match.group(1).strip()
                valid, error = _validate_table(fix_code)
                if valid:
                    return {"convertible": True, "type": "table", "description": desc,
                            "code": fix_code, "route": f"table_fix_{attempt}"}

        return {"convertible": False, "type": "table", "description": desc,
                "code": "", "route": "table_failed"}

    else:
        # code / list 等类型，不做严格校验，直接返回
        return {"convertible": True, "type": content_type, "description": desc,
                "code": code, "route": f"{content_type}_ok"}


def process_file_pipeline(
    file_path: str,
    llm_url: str,
    llm_api_key: str,
    llm_model: str,
    source_dir: Optional[str] = None,
) -> Dict[str, Any]:
    """完整的文件处理流水线。

    流程:
    1. 读取文件（当前仅支持 .md）
    2. 解析 Markdown → 提取文本 + 图片占位符
    3. 复制图片和 md 到临时目录
    4. Phase 1: Prefix cache 描述 — 文档全文放入 system message，全并发 (prefix cache 在 system 前缀中)
    5. Phase 2: 并发判断所有图片可转换性 + 生成代码 (校验 + 重试)
    6. 生成增强 Markdown（含生成代码）写入临时目录
    7. 输出结果

    整个过程不修改源文件，全部在临时目录操作。
    全部使用流式输出 + 思考模式。

    Returns dict:
    - ok, temp_dir, md_file, total_images
    - descriptions: {img_name: description}
    - conversions: [{img_name, type, code, description, route}]
    - output_md: 增强后的 markdown 文件路径
    - errors, type_counts

    Args:
        file_path: MD 文件路径（非 MD 文件应先通过 convert_file_to_markdown 转换）
        llm_url / llm_api_key / llm_model: LLM 配置
        source_dir: 可选，图片所在源目录（由 convert_file_to_markdown 创建的临时目录）。
                   如果不提供，则从 file_path 的同级目录查找图片。
        temp_dir: 可选，与 source_dir 共享的临时目录。如果不提供则创建新的。
    """
    _L = "[pipeline]"
    log = lambda msg: _md_log(f"{_L} {msg}")

    result: Dict[str, Any] = {
        "ok": True,
        "file_path": file_path,
        "temp_dir": "",
        "md_file": "",
        "total_images": 0,
        "descriptions": {},
        "conversions": [],
        "output_md": "",
        "errors": 0,
        "type_counts": {"mermaid": 0, "table": 0, "code": 0, "list": 0, "not_convertible": 0},
    }

    log(f"====== START: {file_path}")

    if not os.path.isfile(file_path):
        log(f"ERROR: file not found: {file_path}")
        result["ok"] = False
        result["error"] = f"File not found: {file_path}"
        return result

    # ---- Step 1: 创建/使用临时目录 ----
    if source_dir and os.path.isdir(source_dir):
        # 复用 convert_file_to_markdown 创建的临时目录
        temp_dir = source_dir
        _owns_temp_dir = False
    else:
        temp_dir = tempfile.mkdtemp(prefix="file_pipeline_")
        _owns_temp_dir = True
    result["temp_dir"] = temp_dir
    log(f"Temp dir: {temp_dir} (owns={_owns_temp_dir})")

    file_ext = os.path.splitext(file_path)[1].lower()
    # MD 文件保持原文件名；非 MD 文件用转换器命名规则 (如 1.docx → 1_docx.md)
    md_filename = os.path.basename(file_path) if file_ext == ".md" else _get_md_output_name(file_path)

    temp_md_path = os.path.join(temp_dir, md_filename)
    # 如果复用 source_dir，转换器已创建 {original_name}.md，不需要复制
    if not os.path.isfile(temp_md_path):
        shutil.copy2(file_path, temp_md_path)

    # 读取 markdown 内容
    try:
        with open(temp_md_path, "r", encoding="utf-8") as fh:
            md_content = fh.read()
    except Exception as exc:
        log(f"ERROR reading file: {exc}")
        result["ok"] = False
        result["error"] = f"Failed to read: {exc}"
        return result

    # 图片查找目录：优先 source_dir（转换器提取的图片），否则用 file_path 同级
    md_dir = source_dir if source_dir else os.path.dirname(os.path.abspath(file_path))
    log(f"Source dir: {md_dir}")
    log(f"File size: {len(md_content)} chars")

    # ---- Step 2: 提取文本 + 图片占位符 ----
    text_with_placeholders, image_infos = _extract_text_and_placeholders(md_content)
    result["total_images"] = len(image_infos)
    log(f"Extracted: {len(image_infos)} image(s)")

    if not image_infos:
        log("No images found, writing copy to temp dir")
        result["output_md"] = temp_md_path
        return result

    for info in image_infos:
        log(f"  [{info['index']}] {info['image_ref']} (line {info['line_number']})")

    # ---- 复制图片到临时目录 + 编码 ----
    image_data = {}  # img_name → {"path", "data_url", "info"}
    for info in image_infos:
        image_ref = info["image_ref"]
        img_name = os.path.basename(image_ref.split("#")[0].split("?")[0])

        # 跳过远程 URL
        if image_ref.startswith(("http://", "https://", "data:")):
            log(f"  [{info['index']}] SKIP: remote URL")
            continue

        # 解析图片路径
        image_path = _resolve_image_path(image_ref, md_dir)
        if image_path is None:
            log(f"  [{info['index']}] SKIP: not found ({image_ref})")
            result["errors"] += 1
            continue

        # 复制到临时目录
        temp_img_path = os.path.join(temp_dir, img_name)
        if not os.path.exists(temp_img_path):
            shutil.copy2(image_path, temp_img_path)

        # 编码为 data URL
        data_url = _encode_image_to_data_url(image_path)
        if not data_url:
            log(f"  [{info['index']}] SKIP: encode failed ({img_name})")
            result["errors"] += 1
            continue

        image_data[img_name] = {
            "path": temp_img_path,
            "data_url": data_url,
            "info": info,
        }

    processable = list(image_data.items())
    log(f"Processable images: {len(processable)}")

    if not processable:
        result["output_md"] = temp_md_path
        return result

    # ---- Phase 1: Prefix cache 描述 (全并发) ----
    # 文档全文已放入 system message，所有请求的 system 前缀完全相同，
    # API 侧 prefix cache 可命中 system + 全文（占总 token 95%+），全并发即可生效
    _pipeline_start = time.time()
    log(f"Phase 1: 描述图片 (全并发, prefix cache 在 system message 中)...")
    _phase1_start = time.time()
    descriptions = {}

    def _describe_task(name, data):
        info = data["info"]
        prompt = DESCRIBE_IMAGE_USER_PROMPT_TEMPLATE.format(idx=info["index"])
        d = _call_vision_llm_with_context(
            data["data_url"], text_with_placeholders, prompt,
            llm_url, llm_api_key, llm_model,
        )
        return name, d

    log(f"  并发描述 {len(processable)} 张图片...")
    with ThreadPoolExecutor(max_workers=min(8, len(processable))) as executor:
        futures = {
            executor.submit(_describe_task, name, data): name
            for name, data in processable
        }
        for future in as_completed(futures):
            name = futures[future]
            try:
                _, desc = future.result()
                descriptions[name] = desc or "(描述失败)"
                log(f"  → {name}: {(desc or '')[:80]}")
            except Exception as e:
                log(f"  → {name}: ERROR {e}")
                descriptions[name] = "(描述失败)"

    result["descriptions"] = descriptions
    _phase1_elapsed = time.time() - _phase1_start
    log(f"Phase 1 完成: {_phase1_elapsed:.1f}s (描述 {len(descriptions)} 张图片)")

    # ---- Phase 2: 并发判断可转换性 + 生成代码 ----
    log(f"Phase 2: 判断转换 + 生成代码 (并发)...")
    _phase2_start = time.time()
    conversions = {}

    def _convert_task(name, data):
        desc = descriptions.get(name, "")
        # 预筛：合并单元格表格直接跳过
        if _is_unconvertible_table(desc):
            print(f"[pipeline:convert] {name}: 预筛命中（不可转换表格），跳过 LLM 调用", flush=True)
            return name, {
                "convertible": False, "type": "table",
                "description": desc, "code": "",
                "route": "table_skipped",
            }
        r = _judge_and_convert_image(
            data["data_url"], desc,
            llm_url, llm_api_key, llm_model,
        )
        return name, r

    with ThreadPoolExecutor(max_workers=10) as executor:
        futures = {
            executor.submit(_convert_task, name, data): name
            for name, data in processable
        }
        for future in as_completed(futures):
            name = futures[future]
            try:
                _, conv = future.result()
                conversions[name] = conv
                ctype = conv.get("type", "")
                route = conv.get("route", "")
                log(f"  → {name}: type={ctype}, route={route}")
            except Exception as e:
                log(f"  → {name}: ERROR {e}")
                conversions[name] = {
                    "convertible": False, "type": "",
                    "description": descriptions.get(name, ""),
                    "code": "", "route": "error",
                }
                result["errors"] += 1

    # 统计
    type_counts = result["type_counts"]
    for name, conv in conversions.items():
        if conv.get("convertible"):
            t = conv.get("type", "")
            if t in type_counts:
                type_counts[t] += 1
        else:
            type_counts["not_convertible"] += 1

    result["conversions"] = [
        {"img_name": name, **conv} for name, conv in conversions.items()
    ]

    _phase2_elapsed = time.time() - _phase2_start
    log(f"Phase 2 完成: {_phase2_elapsed:.1f}s (转换 {len(conversions)} 张图片, types={result['type_counts']})")

    # ---- Step 3: 写入增强 Markdown 到临时目录 ----
    log(f"Phase 3: 写入增强 Markdown（描述插入到图片位置）...")
    _phase3_start = time.time()

    # 将每个占位符替换为：原图 + 描述 + 转换结果（内联插入）
    enhanced_text = text_with_placeholders

    for name, data in processable:
        info = data["info"]
        placeholder = info["placeholder"]
        conv = conversions.get(name, {})
        desc = descriptions.get(name, "")
        ctype = conv.get("type", "")
        code = conv.get("code", "")

        # 构建内联替换内容
        inline_parts = []

        # 原始图片链接
        inline_parts.append(f"![{info['alt_text']}]({name})")

        # 图片描述（用 blockquote 样式）
        if desc:
            inline_parts.append("")
            inline_parts.append(f"> **图片描述**: {desc}")

        # 转换结果
        if conv.get("convertible") and code:
            inline_parts.append("")
            if ctype == "mermaid":
                inline_parts.append("```mermaid")
                inline_parts.append(code)
                inline_parts.append("```")
            elif ctype == "table":
                inline_parts.append(code)
            elif ctype == "code":
                if code.startswith("```"):
                    inline_parts.append(code)
                else:
                    inline_parts.append("```")
                    inline_parts.append(code)
                    inline_parts.append("```")
            elif ctype == "list":
                inline_parts.append(code)

        # 用换行连接，替换占位符
        replacement = "\n".join(inline_parts)
        enhanced_text = enhanced_text.replace(placeholder, replacement)

    # 最终输出
    output_content = f"<!-- Generated by file_pipeline -->\n<!-- Source: {file_path} -->\n\n{enhanced_text}"
    # 输出文件名：MD 文件保持原名，非 MD 文件用带扩展名后缀的命名
    out_name = md_filename  # 复用上面已计算好的文件名
    output_md_path = os.path.join(temp_dir, out_name)
    try:
        with open(output_md_path, "w", encoding="utf-8") as fh:
            fh.write(output_content)
        result["output_md"] = output_md_path
        log(f"Output written: {output_md_path}")
    except Exception as exc:
        log(f"ERROR writing output: {exc}")
        result["ok"] = False
        result["error"] = f"Failed to write output: {exc}"

    # 收集临时目录中的图片文件列表（供导入端复制到 KB 的 raw/assets/）
    try:
        result["image_files"] = [
            f for f in os.listdir(temp_dir)
            if os.path.splitext(f)[1].lower() in IMAGE_EXTENSIONS
        ]
    except Exception:
        result["image_files"] = []

    _phase3_elapsed = time.time() - _phase3_start
    _total_elapsed = time.time() - _pipeline_start
    log(f"Phase 3 完成: {_phase3_elapsed:.1f}s")
    log(
        f"====== DONE: total={result['total_images']}, "
        f"converted={sum(1 for c in conversions.values() if c.get('convertible'))}, "
        f"errors={result['errors']}, types={type_counts}"
    )
    log(
        f"====== 耗时统计: Phase1={_phase1_elapsed:.1f}s + Phase2={_phase2_elapsed:.1f}s + Phase3={_phase3_elapsed:.1f}s = 总计 {_total_elapsed:.1f}s"
    )
    return result


# ============================================================================
#  File → Markdown converters
#  All non-MD files are first converted to Markdown, then fed into the pipeline
# ============================================================================

def _get_md_output_name(source_path: str) -> str:
    """根据源文件名生成对应的 Markdown 输出文件名（保留原始名称 + 原始扩展名，避免同名冲突）。

    例如: '报告.pdf' → '报告_pdf.md', 'slides.pptx' → 'slides_pptx.md'
    这样 1.md 和 1.docx 分别输出为 1.md 和 1_docx.md，不会互相覆盖。
    """
    basename = os.path.basename(source_path)
    name, ext = os.path.splitext(basename)
    ext = ext.lstrip(".").lower()
    if ext:
        return f"{name}_{ext}.md"
    return f"{name}.md"


def _safe_ascii_stem(source_path: str) -> str:
    """从源文件名提取 ASCII 安全前缀（过滤掉非 ASCII 字符如中文；若结果为空则回退为 'img'）。

    例如: 'Markdown渲染测试文件.pdf' → 'Markdown', '渲染测试文件.pdf' → 'img'
    用于图片文件名前缀，确保 100% ASCII 安全，无需 URL 编码。
    """
    basename = os.path.basename(source_path)
    name = os.path.splitext(basename)[0]
    ascii_part = "".join(c for c in name if c.isascii() and (c.isalnum() or c in "_-"))
    return ascii_part if ascii_part else "img"


def _gen_ascii_img_name(stem: str, index: int, ext: str = ".png") -> str:
    """生成纯 ASCII 图片文件名: {stem}_{index}_{随机6位}{ext}"""
    import secrets
    rand = secrets.token_hex(3)  # 6 位十六进制
    if ext == "jpeg":
        ext = "jpg"
    if not ext.startswith("."):
        ext = f".{ext}"
    return f"{stem}_{index}_{rand}{ext}"


def _extract_images_from_pdf(source_path: str, temp_dir: str) -> list:
    """从 PDF 提取图片到 temp_dir，返回图片文件名列表。"""
    try:
        from pypdf import PdfReader
    except ImportError:
        print("[file_convert] pypdf not installed, cannot extract PDF images", flush=True)
        return []

    try:
        reader = PdfReader(source_path)
    except Exception as e:
        print(f"[file_convert] PDF open error: {e}", flush=True)
        return []

    stem = _safe_ascii_stem(source_path)
    img_names = []
    img_count = 0

    for page_num, page in enumerate(reader.pages):
        try:
            for img_obj in page.images:
                img_count += 1
                orig_ext = os.path.splitext(img_obj.name or ".png")[1] or ".png"
                img_name = _gen_ascii_img_name(stem, img_count, orig_ext)

                img_path = os.path.join(temp_dir, img_name)
                try:
                    with open(img_path, "wb") as f:
                        f.write(img_obj.data)
                    img_names.append(img_name)
                except Exception as e:
                    print(f"[file_convert] PDF image save error: {e}", flush=True)
        except Exception as e:
            print(f"[file_convert] PDF image extract error (page {page_num + 1}): {e}", flush=True)

    return img_names


def _extract_images_from_docx(source_path: str, temp_dir: str) -> list:
    """从 DOCX 提取所有嵌入图片到 temp_dir，返回图片文件名列表。

    遍历文档正文及页眉/页脚部件的关系，提取 image 类型部件，
    覆盖 inline/floating/header/footer 等所有图片形式。
    """
    try:
        from docx import Document
    except ImportError:
        print("[file_convert] python-docx not installed, cannot extract DOCX images", flush=True)
        return []

    try:
        doc = Document(source_path)
    except Exception as e:
        print(f"[file_convert] DOCX open error: {e}", flush=True)
        return []

    stem = _safe_ascii_stem(source_path)
    img_names = []
    img_count = 0

    # 收集需要检查的 parts：正文 + 页眉/页脚
    parts_to_check = [doc.part]
    try:
        for section in doc.sections:
            for hf in (section.header, section.footer,
                       section.first_page_header, section.first_page_footer,
                       section.even_page_header, section.even_page_footer):
                try:
                    if hf is not None and hf.part is not None:
                        parts_to_check.append(hf.part)
                except Exception:
                    pass
    except Exception:
        pass

    # 去重（同一 part 可能被多个 section 引用）
    seen = set()
    unique_parts = []
    for p in parts_to_check:
        pid = id(p)
        if pid not in seen:
            seen.add(pid)
            unique_parts.append(p)

    # 遍历所有 parts 的关系，提取 image 类型部件
    try:
        for part in unique_parts:
            for rel_id, rel in part.rels.items():
                if "image" not in rel.reltype:
                    continue
                try:
                    img_count += 1
                    img_ext = rel.target_part.content_type.split("/")[-1].split(";")[0]
                    img_name = _gen_ascii_img_name(stem, img_count, f".{img_ext}")
                    img_path = os.path.join(temp_dir, img_name)
                    with open(img_path, "wb") as f:
                        f.write(rel.target_part.blob)
                    img_names.append(img_name)
                except Exception:
                    continue
    except Exception as e:
        print(f"[file_convert] DOCX image extract error: {e}", flush=True)

    return img_names


def _extract_images_from_pptx(source_path: str, temp_dir: str) -> list:
    """从 PPTX 提取图片到 temp_dir，返回图片文件名列表。"""
    try:
        from pptx import Presentation
    except ImportError:
        print("[file_convert] python-pptx not installed, cannot extract PPTX images", flush=True)
        return []

    try:
        prs = Presentation(source_path)
    except Exception as e:
        print(f"[file_convert] PPTX open error: {e}", flush=True)
        return []

    stem = _safe_ascii_stem(source_path)
    img_names = []
    img_count = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    img_count += 1
                    img_ext = shape.image.content_type.split("/")[-1].split(";")[0]
                    img_name = _gen_ascii_img_name(stem, img_count, f".{img_ext}")
                    img_path = os.path.join(temp_dir, img_name)
                    with open(img_path, "wb") as f:
                        f.write(shape.image.blob)
                    img_names.append(img_name)
                except Exception:
                    pass

    return img_names


def _convert_image_file_to_markdown(source_path: str, temp_dir: str) -> Optional[str]:
    """将单张图片文件转为 output.md（图片引用）。"""
    stem = _safe_ascii_stem(source_path)
    orig_ext = os.path.splitext(source_path)[1].lower() or ".png"
    img_name = _gen_ascii_img_name(stem, 1, orig_ext)

    dest_path = os.path.join(temp_dir, img_name)
    try:
        shutil.copy2(source_path, dest_path)
    except Exception as e:
        print(f"[file_convert] image copy error: {e}", flush=True)
        return None

    caption = stem
    md_content = f"# {caption}\n\n![{caption}]({img_name})\n"
    output_path = os.path.join(temp_dir, _get_md_output_name(source_path))
    try:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(md_content)
        return output_path
    except Exception:
        return None


# 需要图片提取的文件类型 → 图片提取函数
_IMAGE_EXTRACTION_EXTENSIONS = {
    ".pdf": _extract_images_from_pdf,
    ".docx": _extract_images_from_docx,
    ".pptx": _extract_images_from_pptx,
}

# 单张图片文件扩展名
_IMAGE_FILE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"}

# MarkItDown 懒加载单例（避免批量导入时重复加载 onnxruntime 模型）
_markitdown_instance = None


def _get_markitdown():
    """返回模块级 MarkItDown 单例，首次调用时加载 magika/onnxruntime 模型。"""
    global _markitdown_instance
    if _markitdown_instance is None:
        from markitdown import MarkItDown
        _markitdown_instance = MarkItDown()
    return _markitdown_instance


def convert_file_to_markdown(source_path: str) -> Optional[str]:
    """将任意支持的文件类型转换为 Markdown。

    - Markdown 文件：直接返回原路径（流水线会自己处理图片）
    - 单张图片文件：直接生成图片引用 MD
    - 其他文件：markitdown 转文本 + 图片提取（PDF/DOCX/PPTX）→ 合并 MD

    返回 output.md 的路径，或 None（转换失败）。
    调用方负责清理返回的临时目录（可通过 os.path.dirname(result) 获取）。
    """
    ext = os.path.splitext(source_path)[1].lower()
    basename = os.path.basename(source_path)

    # Markdown 文件：不做转换
    if ext in {".md", ".mdx", ".markdown"}:
        return source_path

    # 单张图片文件：直接生成图片引用 MD
    if ext in _IMAGE_FILE_EXTENSIONS:
        temp_dir = tempfile.mkdtemp(prefix="file_convert_")
        try:
            result_path = _convert_image_file_to_markdown(source_path, temp_dir)
            if result_path and os.path.isfile(result_path):
                print(f"[file_convert] {basename} -> {result_path}", flush=True)
                return result_path
            else:
                shutil.rmtree(temp_dir, ignore_errors=True)
                return None
        except Exception as e:
            print(f"[file_convert] {basename} image convert error: {e}", flush=True)
            shutil.rmtree(temp_dir, ignore_errors=True)
            return None

    # 其他文件：markitdown 转文本 + 图片提取
    temp_dir = tempfile.mkdtemp(prefix="file_convert_")
    try:
        # Step 1: markitdown 转文本
        md_converter = _get_markitdown()
        result = md_converter.convert(source_path)
        text_content = result.text_content or ""

        if not text_content.strip():
            print(f"[file_convert] {basename} markitdown output empty", flush=True)
            shutil.rmtree(temp_dir, ignore_errors=True)
            return None

        # 清理 markitdown 输出中的控制字符（PDF 提取常含 \x00/\x0c 等）
        # 保留 \t \n \r，移除其他 C0 控制字符
        text_content = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text_content)

        # 清理 markitdown 输出中的 base64 data URI 图片引用
        # markitdown 会将 DOCX inline 图片输出为 ![name](data:image/png;base64,...)
        # 我们单独提取图片，所以移除这些 data URI 引用
        text_content = re.sub(r'!\[[^\]]*\]\(data:image/[^)]+\)', '', text_content)

        # Step 2: 图片提取（如有）
        image_lines = []
        image_extractor = _IMAGE_EXTRACTION_EXTENSIONS.get(ext)
        if image_extractor:
            img_names = image_extractor(source_path, temp_dir)
            for img_name in img_names:
                image_lines.append(f"![{img_name}]({img_name})")

        # Step 3: 合并文本 + 图片引用
        if image_lines:
            md_content = text_content + "\n\n---\n\n" + "\n".join(image_lines) + "\n"
        else:
            md_content = text_content

        # Step 4: 写入 MD 文件
        output_path = os.path.join(temp_dir, _get_md_output_name(source_path))
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(md_content)

        print(f"[file_convert] {basename} -> {output_path} (text={len(text_content)} chars, images={len(image_lines)})", flush=True)
        return output_path
    except Exception as e:
        print(f"[file_convert] {basename} convert error: {e}", flush=True)
        shutil.rmtree(temp_dir, ignore_errors=True)
        return None
